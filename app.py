import json, re, os, time
import streamlit as st
import streamlit.components.v1
import io, pandas as pd
from dotenv import load_dotenv
from statistics import mean
from datetime import datetime

load_dotenv()

from openai import OpenAI
from auth import require_auth, get_current_username, get_current_role, logout_user
from test_manager import (
    generate_test_id, validate_test_data, convert_to_rubric_format
)

DEEPSEEK_API_KEY = os.getenv("DEEPSEEK_API_KEY")
client = None
if DEEPSEEK_API_KEY:
    client = OpenAI(api_key=DEEPSEEK_API_KEY, base_url="https://api.deepseek.com")  # OpenAI-compatible
DEEPSEEK_MODEL = "deepseek-chat"  # or "deepseek-reasoner"


# --- persistence + badges
POINTS_PATH = "points.json"
SUBM_PATH   = "submissions.json"
REVIEWS_PATH = "reviews.json"

# --- Materials storage
MATERIALS_DIR = "materials"
MATERIALS_INDEX = "materials_index.json"  # subject -> [ {name, path, size, mime, uploader, note, ts} ]


def current_subject() -> str:
    return st.session_state.get("subject", "General")

def assignment_selector(widget_key: str):
    subj = current_subject()
    sub_rubrics = [r for r in RUBRICS if r.get("subject", "General") == subj]
    if not sub_rubrics:
        st.info("Для этого предмета пока нет заданий.")
        return None
    id_to_r = {r["assignment_id"]: r for r in sub_rubrics}
    ids = list(id_to_r.keys())

    cur = st.session_state.get("assign_id", ids[0])
    if cur not in ids:
        cur = ids[0]

    def format_assignment(rid: str) -> str:
        r = id_to_r[rid]
        title = r.get("title", rid)
        due = r.get("due_date")
        if due:
            try:
                # Показать как ДД.ММ.ГГГГ
                from datetime import datetime
                d = datetime.fromisoformat(due)
                return f"{title} (до {d.strftime('%d.%m.%Y')})"
            except Exception:
                return f"{title} (до {due})"
        return title

    rid = st.selectbox(
        "Задание:",
        options=ids,
        index=ids.index(cur),
        format_func=format_assignment,
        key=widget_key,
    )
    st.session_state["assign_id"] = rid
    return id_to_r[rid]

def submission_preview(s: dict, max_len: int = 60) -> str:
    """Короткий текст для списка выбора в кросс-проверке.
    Поддерживает старые ('answer') и новые ('answers'=[...]) форматы."""
    if "answer" in s:  
        txt = s.get("answer") or ""
    elif "answers" in s and s["answers"]:
        # возьмём первый ответ, либо склеим первые строки
        parts = [a.get("answer", "") for a in s["answers"] if a.get("answer")]
        txt = " | ".join(parts[:2]) if parts else ""
    else:
        txt = ""
    txt = txt.strip().replace("\n", " ")
    return (txt[:max_len] + "…") if len(txt) > max_len else txt


def load_json(path, default):
    try:
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return default

def save_json(path, obj):
    try:
        with open(path, "w", encoding="utf-8") as f:
            json.dump(obj, f, ensure_ascii=False, indent=2)
    except Exception:
        pass

def badge_for(points: int) -> str:
    if points >= 50: return "🏅 Power Learner"
    if points >= 30: return "🥇 Consistent"
    if points >= 15: return "🥈 Getting There"
    if points >= 5:  return "🥉 Starter"
    return "🌱 Rookie"

def migrate_stores_to_subject_scope():
    # был dict user->points, теперь dict subject->(dict user->points)
    if isinstance(st.session_state.points, dict) and st.session_state.points and \
       all(isinstance(v, int) for v in st.session_state.points.values()):
        st.session_state.points = {"General": st.session_state.points}
        save_json(POINTS_PATH, st.session_state.points)

    # dict subject->list
    if isinstance(st.session_state.submissions, list):
        st.session_state.submissions = {"General": st.session_state.submissions}
        save_json(SUBM_PATH, st.session_state.submissions)

    # dict subject->list
    if isinstance(st.session_state.reviews, list):
        st.session_state.reviews = {"General": st.session_state.reviews}
        save_json(REVIEWS_PATH, st.session_state.reviews)

def load_materials_index():
    return load_json(MATERIALS_INDEX, {})

def save_materials_index(idx):
    save_json(MATERIALS_INDEX, idx)

def ensure_dir(path: str):
    os.makedirs(path, exist_ok=True)

def get_materials_store():
    idx = load_materials_index()
    subj = current_subject()
    if subj not in idx:
        idx[subj] = []
        save_materials_index(idx)
    return idx, idx[subj]

def safe_filename(name: str) -> str:
    """Создает безопасное имя файла для файловой системы"""
    return re.sub(r"[^A-Za-z0-9._\-]+", "_", name)

def display_filename(name: str, original_name: str = None) -> str:
    """Возвращает читаемое название файла для отображения"""
    # Если есть оригинальное название, используем его
    if original_name and original_name.strip():
        return original_name
    
    # Если имя содержит только подчеркивания и расширение, показываем более понятное название
    clean_name = name.replace("_", "").replace(".", "").replace("-", "")
    if len(clean_name) <= 2:  # Только расширение или очень короткое название
        # Пытаемся определить тип файла по расширению
        if name.lower().endswith('.pdf'):
            return "Документ PDF"
        elif name.lower().endswith('.docx'):
            return "Документ Word"
        elif name.lower().endswith('.pptx'):
            return "Презентация PowerPoint"
        elif name.lower().endswith('.txt'):
            return "Текстовый файл"
        else:
            return "Файл"
    
    # Если название состоит в основном из подчеркиваний, но есть некоторые символы
    if name.count('_') > len(name) * 0.7:  # Больше 70% подчеркиваний
        # Пытаемся восстановить читаемое название
        parts = name.split('_')
        readable_parts = []
        for part in parts:
            if part and not part.startswith('.'):
                readable_parts.append(part)
        
        if readable_parts:
            return ' '.join(readable_parts)
        else:
            return "Файл"
    
    return name

def extract_text_from_file(file_path: str, mime_type: str) -> str:
    """Извлекает текст из файла разных форматов"""
    try:
        if mime_type == "application/pdf" or file_path.lower().endswith('.pdf'):
            try:
                import PyPDF2
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page in pdf_reader.pages:
                        page_text = page.extract_text()
                        # Исправляем проблемы с кодировкой
                        if page_text:
                            # Пытаемся исправить кодировку для русских символов
                            try:
                                # Если текст содержит неправильно декодированные символы
                                if 'Ð' in page_text or 'Ñ' in page_text:
                                    # Пытаемся перекодировать
                                    page_text = page_text.encode('latin1').decode('utf-8', errors='ignore')
                            except:
                                pass
                            text += page_text + "\n"
                    return text.strip()
            except Exception as pdf_error:
                # Fallback: пытаемся использовать pdfplumber если доступен
                try:
                    import pdfplumber
                    with pdfplumber.open(file_path) as pdf:
                        text = ""
                        for page in pdf.pages:
                            page_text = page.extract_text()
                            if page_text:
                                text += page_text + "\n"
                        return text.strip()
                except ImportError:
                    return f"Ошибка извлечения текста из PDF: {str(pdf_error)}. Установите pdfplumber для лучшей поддержки кодировки."
        
        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or file_path.lower().endswith('.docx'):
            try:
                from docx import Document
                doc = Document(file_path)
                text = ""
                for paragraph in doc.paragraphs:
                    if paragraph.text:
                        text += paragraph.text + "\n"
                return text.strip()
            except Exception as docx_error:
                return f"Ошибка извлечения текста из DOCX: {str(docx_error)}"
        
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or file_path.lower().endswith('.pptx'):
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text += shape.text + "\n"
                return text.strip()
            except Exception as pptx_error:
                return f"Ошибка извлечения текста из PPTX: {str(pptx_error)}"
        
        elif mime_type == "text/plain" or file_path.lower().endswith('.txt'):
            # Пытаемся разные кодировки для TXT файлов
            encodings = ['utf-8', 'cp1251', 'latin1', 'utf-16']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        return file.read().strip()
                except UnicodeDecodeError:
                    continue
            return "Ошибка: не удалось определить кодировку текстового файла"
        
        else:
            return f"Формат файла {mime_type} не поддерживается для извлечения текста"
    
    except Exception as e:
        return f"Ошибка извлечения текста: {str(e)}"

def generate_annotation(text: str, filename: str) -> str:
    """Генерирует аннотацию для текста через DeepSeek API"""
    if client is None:
        return "AI недоступен: не найден ключ API"
    
    if not text or len(text.strip()) < 50:
        return "Недостаточно текста для генерации аннотации"
    
    try:
        # Ограничиваем длину текста для API (примерно 4000 символов)
        max_length = 4000
        if len(text) > max_length:
            text = text[:max_length] + "..."
        
        system_msg = (
            "Ты - эксперт по созданию кратких аннотаций учебных материалов. "
            "Создай краткую аннотацию (2-4 предложения) на русском языке, которая: "
            "1) Описывает основную тему материала "
            "2) Выделяет ключевые понятия и идеи "
            "3) Указывает на практическую ценность "
            "4) Пишется понятным языком для студентов"
        )
        
        user_msg = f"Файл: {filename}\n\nТекст материала:\n{text}"
        
        resp = client.chat.completions.create(
            model=DEEPSEEK_MODEL,
            temperature=0.3,
            messages=[
                {"role": "system", "content": system_msg},
                {"role": "user", "content": user_msg}
            ],
        )
        
        annotation = resp.choices[0].message.content.strip()
        return annotation
    
    except Exception as e:
        return f"Ошибка генерации аннотации: {str(e)}"

def get_points_store() -> dict:
    subj = current_subject()
    if subj not in st.session_state.points:
        st.session_state.points[subj] = {}
    return st.session_state.points[subj]

def save_points_store():
    save_json(POINTS_PATH, st.session_state.points)

def get_submissions_store() -> list:
    subj = current_subject()
    if subj not in st.session_state.submissions:
        st.session_state.submissions[subj] = []
    return st.session_state.submissions[subj]

def save_submissions_store():
    save_json(SUBM_PATH, st.session_state.submissions)

def get_reviews_store() -> list:
    subj = current_subject()
    if subj not in st.session_state.reviews:
        st.session_state.reviews[subj] = []
    return st.session_state.reviews[subj]

def save_reviews_store():
    save_json(REVIEWS_PATH, st.session_state.reviews)

def process_chat_question(question: str, subject: str) -> str:
    """
    Обрабатывает вопрос пользователя через нейронную сеть.
    Анализирует дедлайны, содержание материалов и дает рекомендации.
    """
    if client is None:
        return "❌ AI недоступен: не найден ключ API (DEEPSEEK_API_KEY)"
    
    try:
        # Собираем контекст для AI
        context = build_chat_context(subject)
        
        # Определяем тип вопроса и формируем ответ
        question_lower = question.lower()
        
        if any(word in question_lower for word in ["дедлайн", "срок", "когда", "до какого"]):
            return handle_deadline_question(question, subject, context)
        elif any(word in question_lower for word in ["материал", "лекция", "прочитать", "изучить", "что читать", "задание"]):
            return handle_material_question(question, subject, context)
        elif any(word in question_lower for word in ["просрочен", "опоздал", "не успел"]):
            return handle_overdue_question(question, subject, context)
        else:
            return handle_general_question(question, subject, context)
            
    except Exception as e:
        return f"❌ Ошибка обработки вопроса: {str(e)}"

def build_chat_context(subject: str) -> dict:
    """Собирает контекстную информацию для AI"""
    context = {
        "subject": subject,
        "assignments": [],
        "materials": [],
        "videos": [],
        "current_date": datetime.now().isoformat()
    }
    
    # Собираем информацию об заданиях
    subj_rubrics = [r for r in RUBRICS if r.get("subject", "General") == subject]
    for rubric in subj_rubrics:
        assignment_info = {
            "title": rubric.get("title", ""),
            "assignment_id": rubric.get("assignment_id", ""),
            "due_date": rubric.get("due_date", ""),
            "questions": []
        }
        
        # Добавляем информацию о вопросах
        for q in rubric.get("questions", []):
            keywords = []
            for kw in q.get("keywords", []):
                if isinstance(kw, dict):
                    keywords.append(kw.get("word", ""))
                else:
                    keywords.append(str(kw))
            
            question_info = {
                "title": q.get("title", ""),
                "keywords": keywords
            }
            assignment_info["questions"].append(question_info)
        
        context["assignments"].append(assignment_info)
    
    # Собираем информацию о материалах
    idx, materials = get_materials_store()
    for material in materials:
        if isinstance(material, dict):
            material_info = {
                "name": material.get("name", ""),
                "note": material.get("note", ""),
                "annotation": material.get("annotation", "")
            }
            context["materials"].append(material_info)
    
    # Собираем информацию о видео
    try:
        with open("videos.json", "r", encoding="utf-8") as f:
            videos_data = json.load(f)
        subject_videos = videos_data.get(subject, [])
        for video in subject_videos:
            if isinstance(video, dict):
                video_info = {
                    "title": video.get("title", ""),
                    "note": video.get("note", "")
                }
                context["videos"].append(video_info)
    except:
        pass
    
    return context

def handle_deadline_question(question: str, subject: str, context: dict) -> str:
    """Обрабатывает вопросы о дедлайнах"""
    response_parts = []
    
    # Анализируем дедлайны
    today = datetime.now()
    overdue_assignments = []
    upcoming_assignments = []
    
    for assignment in context["assignments"]:
        due_date_str = assignment.get("due_date", "")
        if due_date_str:
            try:
                due_date = datetime.fromisoformat(due_date_str)
                days_left = (due_date - today).days
                
                if days_left < 0:
                    overdue_assignments.append({
                        "title": assignment["title"],
                        "days_overdue": abs(days_left)
                    })
                elif days_left <= 7:  # Ближайшие 7 дней
                    upcoming_assignments.append({
                        "title": assignment["title"],
                        "days_left": days_left,
                        "due_date": due_date_str
                    })
            except:
                pass
    
    if overdue_assignments:
        response_parts.append("🔴 **Просроченные задания:**")
        for assignment in overdue_assignments:
            response_parts.append(f"• **{assignment['title']}** - просрочено на {assignment['days_overdue']} дн.")
    
    if upcoming_assignments:
        response_parts.append("\n🟡 **Ближайшие дедлайны:**")
        for assignment in upcoming_assignments:
            if assignment['days_left'] == 0:
                response_parts.append(f"• **{assignment['title']}** - сегодня!")
            else:
                response_parts.append(f"• **{assignment['title']}** - через {assignment['days_left']} дн.")
    
    if not overdue_assignments and not upcoming_assignments:
        response_parts.append("✅ У вас нет срочных дедлайнов в ближайшие дни.")
    
    return "\n".join(response_parts)

def handle_material_question(question: str, subject: str, context: dict) -> str:
    """Обрабатывает вопросы о материалах и рекомендациях"""
    if client is None:
        return "❌ AI недоступен для анализа материалов"
    
    try:
        # Проверяем, спрашивают ли о конкретном тесте/задании
        question_lower = question.lower()
        specific_assignment = find_specific_assignment(question, context["assignments"])
        
        if specific_assignment:
            return handle_specific_assignment_question(question, specific_assignment, context)
        else:
            # Проверяем, спрашивали ли о конкретном задании
            import re
            numbers = re.findall(r'\d+', question)
            if numbers and any(word in question_lower for word in ["тест", "задание", "дз", "контрольная"]):
                return handle_not_found_assignment(question, numbers, context)
        
        # Анализируем ключевые слова из вопросов для лучших рекомендаций
        relevant_keywords = extract_keywords_from_question(question)
        matching_assignments = find_matching_assignments(question, context["assignments"])
        
        # Проводим детальный анализ материалов
        material_analysis = analyze_materials_for_question(question, context["materials"], relevant_keywords)
        video_analysis = analyze_videos_for_question(question, context["videos"], relevant_keywords)
        
        # Формируем конкретный ответ на основе анализа
        response_parts = []
        
        if material_analysis["suitable_materials"]:
            response_parts.append("📚 **Подходящие материалы:**")
            for material in material_analysis["suitable_materials"]:
                response_parts.append(f"• **{material['name']}** - {material['reason']}")
        else:
            response_parts.append("📚 **Материалы:** Подходящих материалов не найдено.")
        
        if video_analysis["suitable_videos"]:
            response_parts.append("\n🎥 **Подходящие видео:**")
            for video in video_analysis["suitable_videos"]:
                response_parts.append(f"• **{video['title']}** - {video['reason']}")
        else:
            response_parts.append("\n🎥 **Видео:** Подходящих видео не найдено.")
        
        if matching_assignments:
            response_parts.append("\n📝 **Релевантные задания:**")
            for assignment in matching_assignments:
                response_parts.append(f"• **{assignment['title']}**")
                for q in assignment.get("questions", []):
                    if any(keyword.lower() in question.lower() for keyword in q.get("keywords", [])):
                        response_parts.append(f"  - {q['title']}")
        
        # Если ничего не найдено, даем общий совет
        if not material_analysis["suitable_materials"] and not video_analysis["suitable_videos"]:
            response_parts.append("\n💡 **Рекомендации:**")
            response_parts.append("• Обратитесь к преподавателю за дополнительными материалами")
            response_parts.append("• Используйте внешние источники по теме")
            response_parts.append("• Проверьте, есть ли материалы в других предметах")
        
        return "\n".join(response_parts)
        
    except Exception as e:
        return f"❌ Ошибка анализа материалов: {str(e)}"

def find_specific_assignment(question: str, assignments: list) -> dict:
    """Находит конкретное задание по номеру или названию"""
    question_lower = question.lower()
    
    # Ищем по номеру (тест 32, задание 2, дз1 и т.д.)
    import re
    numbers = re.findall(r'\d+', question)
    
    # Сначала ищем точное совпадение по номеру
    for assignment in assignments:
        assignment_title = assignment.get("title", "").lower().strip()
        assignment_id = assignment.get("assignment_id", "").lower()
        
        # Проверяем точное совпадение по номеру в названии
        for number in numbers:
            # Точное совпадение: номер должен быть равен названию задания
            if number == assignment_title:
                return assignment
            
            # Проверяем, что номер не является частью большего числа
            # Например, "11" не должно находить "1134"
            if number in assignment_title:
                # Проверяем, что это не часть большего числа
                import re
                # Ищем все числа в названии задания
                title_numbers = re.findall(r'\d+', assignment_title)
                if number in title_numbers:
                    # Проверяем, что это точное совпадение, а не часть
                    if number == assignment_title or f" {number} " in f" {assignment_title} " or assignment_title.startswith(number + " ") or assignment_title.endswith(" " + number):
                        return assignment
    
    # Проверяем по ключевым словам (только если есть точное совпадение)
    for assignment in assignments:
        assignment_title = assignment.get("title", "").lower()
        if any(word in assignment_title for word in ["тест", "задание", "дз", "контрольная"]):
            if any(word in question_lower for word in ["тест", "задание", "дз", "контрольная"]):
                # Дополнительная проверка на номер
                for number in numbers:
                    if number in assignment_title:
                        return assignment
    
    return None

def handle_not_found_assignment(question: str, numbers: list, context: dict) -> str:
    """Обрабатывает случай, когда запрашиваемое задание не найдено"""
    response_parts = []
    
    # Показываем, что задание не найдено
    response_parts.append(f"❌ **Задание не найдено**")
    response_parts.append(f"Задание с номером {', '.join(numbers)} не существует в курсе.")
    
    # Показываем доступные задания
    available_assignments = context["assignments"]
    if available_assignments:
        response_parts.append("\n📝 **Доступные задания:**")
        for assignment in available_assignments:
            title = assignment.get("title", "Без названия")
            response_parts.append(f"• **{title}**")
    else:
        response_parts.append("\n📝 **Доступные задания:** Нет активных заданий.")
    
    # Рекомендации
    response_parts.append("\n💡 **Рекомендации:**")
    response_parts.append("• Проверьте правильность номера задания")
    response_parts.append("• Обратитесь к преподавателю для уточнения")
    response_parts.append("• Посмотрите список доступных заданий выше")
    
    return "\n".join(response_parts)

def handle_specific_assignment_question(question: str, assignment: dict, context: dict) -> str:
    """Обрабатывает вопросы о конкретном задании"""
    response_parts = []
    
    # Показываем информацию о задании
    response_parts.append(f"📝 **{assignment['title']}**")
    
    # Анализируем тип задания
    test_type = assignment.get("test_type", "")
    questions = assignment.get("questions", [])
    
    if not questions:
        response_parts.append("\n❌ **В задании нет вопросов.**")
        response_parts.append("Это не полноценный тест - возможно, задание еще не настроено.")
        return "\n".join(response_parts)
    
    # Анализируем качество вопросов
    response_parts.append(f"\n**Анализ теста:**")
    
    # Проверяем, является ли это реальным тестом
    is_real_test = True
    issues = []
    
    for i, q in enumerate(questions, 1):
        question_title = q.get("title", "")
        
        # Проверяем качество вопроса - должен быть осмысленным текстом
        if not question_title or len(question_title.strip()) < 5:
            is_real_test = False
            issues.append(f"Вопрос {i}: отсутствует или слишком короткий ('{question_title}')")
            continue
        
        # Проверяем, что вопрос не является просто числом или бессмысленным текстом
        if question_title.strip().isdigit() or len(question_title.strip()) < 10:
            is_real_test = False
            issues.append(f"Вопрос {i}: не является полноценным вопросом ('{question_title}')")
            continue
        
        # Проверяем варианты ответов для тестов с выбором
        if test_type == "multiple_choice" or q.get("options"):
            options = q.get("options", [])
            if not options or len(options) < 2:
                is_real_test = False
                issues.append(f"Вопрос {i}: недостаточно вариантов ответов")
            elif len(set(options)) < len(options):
                is_real_test = False
                issues.append(f"Вопрос {i}: есть дублирующиеся варианты ответов ({', '.join(options)})")
            elif all(opt.strip().isdigit() and len(opt.strip()) == 1 for opt in options):
                is_real_test = False
                issues.append(f"Вопрос {i}: варианты ответов выглядят как тестовые данные")
        
        # Проверяем ключевые слова для текстовых вопросов
        if test_type == "keyword_based" or q.get("keywords"):
            keywords = q.get("keywords", [])
            if not keywords:
                is_real_test = False
                issues.append(f"Вопрос {i}: нет ключевых слов для оценки")
    
    # Даем честную оценку теста
    if not is_real_test:
        response_parts.append("❌ **Это не полноценный тест.**")
        response_parts.append("Проблемы:")
        for issue in issues:
            response_parts.append(f"• {issue}")
        
        # Проверяем, является ли это тестовыми данными
        if any("тестовые данные" in issue.lower() for issue in issues):
            response_parts.append("\n⚠️ **Это выглядит как тестовые данные для проверки функционала.**")
            response_parts.append("Скорее всего, это не реальное задание для студентов.")
        
        response_parts.append("\n💡 **Рекомендация:** Обратитесь к преподавателю для настройки теста.")
    else:
        response_parts.append("✅ **Это полноценный тест.**")
        response_parts.append(f"Количество вопросов: {len(questions)}")
        response_parts.append(f"Тип теста: {'С вариантами ответов' if test_type == 'multiple_choice' else 'С ключевыми словами'}")
    
    # НЕ показываем детали вопросов - это нарушает принципы обучения
    response_parts.append(f"\n📊 **Информация о тесте:**")
    response_parts.append(f"• Количество вопросов: {len(questions)}")
    response_parts.append(f"• Тип: {'Тест с вариантами ответов' if test_type == 'multiple_choice' else 'Текстовые вопросы'}")
    response_parts.append(f"• Максимальный балл: {sum(q.get('max_points', 0) for q in questions)}")
    
    response_parts.append("\n⚠️ **Важно:** Я не раскрываю конкретные вопросы и ответы - это нарушило бы принципы обучения.")
    
    # Проверяем дедлайн
    due_date = assignment.get("due_date")
    if due_date:
        try:
            from datetime import datetime
            due = datetime.fromisoformat(due_date)
            today = datetime.now()
            days_left = (due - today).days
            
            if days_left < 0:
                response_parts.append(f"\n🔴 **Дедлайн:** Просрочено на {abs(days_left)} дн.")
                response_parts.append("💡 **Рекомендация:** Даже если задание просрочено, его стоит выполнить для изучения материала.")
            elif days_left == 0:
                response_parts.append(f"\n🟡 **Дедлайн:** Сегодня!")
            else:
                response_parts.append(f"\n🟢 **Дедлайн:** Осталось {days_left} дн.")
        except:
            response_parts.append(f"\n📅 **Дедлайн:** {due_date}")
    
    # Анализируем материалы только для полноценных тестов
    if is_real_test:
        # Используем нейронную сеть для анализа материалов
        ai_recommendations = get_ai_material_recommendations(question, assignment, context)
        response_parts.append(ai_recommendations)
    else:
        response_parts.append("\n📚 **Материалы:** Не анализируются для неполноценного теста.")
        response_parts.append("🎥 **Видео:** Не анализируются для неполноценного теста.")
    
    return "\n".join(response_parts)

def generate_ai_test(title: str, subject: str, description: str, question_count: int, selected_materials: list, material_options: dict) -> dict:
    """Генерирует тест с помощью нейронной сети"""
    if client is None:
        return None
    
    try:
        # Собираем информацию о выбранных материалах
        materials_text = ""
        for material_name in selected_materials:
            material_info = material_options[material_name]
            materials_text += f"\n• {material_info['name']}: {material_info.get('note', '')} - {material_info.get('annotation', '')}"
        
        system_msg = (
            "Ты - эксперт по созданию образовательных тестов. "
            "Создай тест с вариантами ответов на основе предоставленных материалов. "
            "Каждый вопрос должен иметь 4 варианта ответа, один из которых правильный. "
            "Вопросы должны проверять понимание материала, а не просто запоминание. "
            "Отвечай строго в JSON формате."
        )
        
        user_msg = f"""
Создай тест со следующими параметрами:

Название: {title}
Предмет: {subject}
Описание: {description}
Количество вопросов: {question_count}

Материалы для генерации вопросов:
{materials_text}

Требования:
1. Создай {question_count} вопросов с вариантами ответов
2. Каждый вопрос должен иметь 4 варианта ответа (A, B, C, D)
3. Один вариант должен быть правильным
4. Вопросы должны быть разного уровня сложности
5. Вопросы должны проверять понимание, а не запоминание
6. Используй информацию из предоставленных материалов

Формат ответа (строго JSON):
{{
  "questions": [
    {{
      "question_id": "q1",
      "title": "Текст вопроса",
      "options": ["Вариант A", "Вариант B", "Вариант C", "Вариант D"],
      "correct_answer": "Вариант A",
      "max_points": 10
    }}
  ]
}}
"""
        
        resp = client.chat.completions.create(
            model=DEEPSEEK_MODEL,
            temperature=0.7,
            messages=[
                {"role": "system", "content": system_msg},
                {"role": "user", "content": user_msg}
            ],
        )
        
        # Парсим JSON ответ
        import json
        content = resp.choices[0].message.content.strip()
        
        # Извлекаем JSON из ответа
        import re
        json_match = re.search(r'\{.*\}', content, re.DOTALL)
        if json_match:
            json_str = json_match.group(0)
            test_data = json.loads(json_str)
            
            # Добавляем метаданные
            test_data["title"] = title
            test_data["subject"] = subject
            test_data["description"] = description
            test_data["test_type"] = "multiple_choice"
            test_data["ai_generated"] = True
            
            return test_data
        else:
            return None
            
    except Exception as e:
        st.error(f"Ошибка генерации теста: {str(e)}")
        return None

def save_generated_test(test_data: dict, time_limit: int, due_date) -> bool:
    """Сохраняет сгенерированный тест в rubric.json"""
    try:
        # Генерируем ID теста
        test_id = generate_test_id(test_data["title"], test_data["subject"])
        test_data["assignment_id"] = test_id
        
        # Добавляем время и дедлайн
        if time_limit > 0:
            test_data["time_limit_minutes"] = time_limit
        
        if due_date:
            test_data["due_date"] = due_date.isoformat()
        
        # Конвертируем в формат рубрики
        rubric_data = convert_to_rubric_format(test_data)
        
        # Загружаем существующие рубрики
        with open("rubric.json", "r", encoding="utf-8") as f:
            rubrics = json.load(f)
        
        # Добавляем новую рубрику
        rubrics.append(rubric_data)
        
        # Сохраняем
        with open("rubric.json", "w", encoding="utf-8") as f:
            json.dump(rubrics, f, ensure_ascii=False, indent=2)
        
        return True
        
    except Exception as e:
        st.error(f"Ошибка сохранения теста: {str(e)}")
        return False

def get_ai_material_recommendations(question: str, assignment: dict, context: dict) -> str:
    """Получает рекомендации материалов через нейронную сеть"""
    if client is None:
        return "\n📚 **Материалы:** AI недоступен для анализа материалов."
    
    try:
        # Формируем информацию о задании для AI
        assignment_title = assignment.get("title", "")
        test_type = assignment.get("test_type", "")
        
        # Собираем информацию о материалах
        materials_info = []
        for material in context["materials"]:
            materials_info.append({
                "name": material.get("name", ""),
                "note": material.get("note", ""),
                "annotation": material.get("annotation", "")
            })
        
        # Собираем информацию о видео
        videos_info = []
        for video in context["videos"]:
            videos_info.append({
                "title": video.get("title", ""),
                "note": video.get("note", "")
            })
        
        system_msg = (
            "Ты - помощник студента по учебным материалам. "
            "Анализируй задание и рекомендуй ТОЛЬКО подходящие материалы. "
            "Будь строгим в отборе - не рекомендуй материалы, которые не связаны с темой. "
            "Отвечай на русском языке, будь конкретным и честным. "
            "Если подходящих материалов нет, честно скажи об этом."
        )
        
        user_msg = f"""
Вопрос студента: {question}

Задание: {assignment_title}
Тип теста: {test_type}

Доступные материалы:
{chr(10).join([f"• {m['name']}: {m.get('note', '')} - {m.get('annotation', '')}" for m in materials_info])}

Доступные видео:
{chr(10).join([f"• {v['title']}: {v.get('note', '')}" for v in videos_info])}

Проанализируй и дай рекомендации ТОЛЬКО по материалам, которые действительно помогут с этим заданием.
Не рекомендуй материалы, которые не связаны с темой (например, музыкальные видео для технических заданий).
"""
        
        resp = client.chat.completions.create(
            model=DEEPSEEK_MODEL,
            temperature=0.2,
            messages=[
                {"role": "system", "content": system_msg},
                {"role": "user", "content": user_msg}
            ],
        )
        
        return "\n🤖 **AI-рекомендации материалов:**\n" + resp.choices[0].message.content.strip()
        
    except Exception as e:
        return f"\n📚 **Материалы:** Ошибка AI-анализа: {str(e)}"

def analyze_materials_for_question(question: str, materials: list, keywords: list) -> dict:
    """Анализирует материалы на предмет соответствия вопросу"""
    question_lower = question.lower()
    suitable_materials = []
    
    # Специальная обработка для вопросов о конкретных заданиях
    if "задание" in question_lower:
        # Ищем материалы, которые могут подходить для любого задания
        for material in materials:
            material_name = material.get("name", "").lower()
            material_note = material.get("note", "").lower()
            material_annotation = material.get("annotation", "").lower()
            
            # Если материал содержит общие термины по ML или Python
            if any(term in material_name or term in material_note or term in material_annotation 
                   for term in ["python", "ml", "машинное", "обучение", "алгоритм", "модель"]):
                suitable_materials.append({
                    "name": material["name"],
                    "reason": "содержит общую информацию по машинному обучению и Python"
                })
    
    # Обычный анализ по ключевым словам
    for material in materials:
        material_name = material.get("name", "").lower()
        material_note = material.get("note", "").lower()
        material_annotation = material.get("annotation", "").lower()
        
        # Проверяем соответствие по названию
        if any(keyword in material_name for keyword in keywords):
            suitable_materials.append({
                "name": material["name"],
                "reason": f"название содержит ключевые слова: {', '.join([k for k in keywords if k in material_name])}"
            })
            continue
        
        # Проверяем соответствие по описанию
        if any(keyword in material_note for keyword in keywords):
            suitable_materials.append({
                "name": material["name"],
                "reason": f"описание содержит ключевые слова: {', '.join([k for k in keywords if k in material_note])}"
            })
            continue
        
        # Проверяем соответствие по аннотации
        if any(keyword in material_annotation for keyword in keywords):
            suitable_materials.append({
                "name": material["name"],
                "reason": f"аннотация содержит ключевые слова: {', '.join([k for k in keywords if k in material_annotation])}"
            })
            continue
        
        # Проверяем общее соответствие по вопросу
        if any(word in material_name or word in material_note or word in material_annotation 
               for word in question_lower.split() if len(word) > 3):
            suitable_materials.append({
                "name": material["name"],
                "reason": "содержит релевантную информацию по теме вопроса"
            })
    
    # Убираем дубликаты
    unique_materials = []
    seen_names = set()
    for material in suitable_materials:
        if material["name"] not in seen_names:
            unique_materials.append(material)
            seen_names.add(material["name"])
    
    return {"suitable_materials": unique_materials}

def analyze_videos_for_question(question: str, videos: list, keywords: list) -> dict:
    """Анализирует видео на предмет соответствия вопросу"""
    question_lower = question.lower()
    suitable_videos = []
    
    for video in videos:
        video_title = video.get("title", "").lower()
        video_note = video.get("note", "").lower()
        
        # Проверяем соответствие по названию
        if any(keyword in video_title for keyword in keywords):
            suitable_videos.append({
                "title": video["title"],
                "reason": f"название содержит ключевые слова: {', '.join([k for k in keywords if k in video_title])}"
            })
            continue
        
        # Проверяем соответствие по описанию
        if any(keyword in video_note for keyword in keywords):
            suitable_videos.append({
                "title": video["title"],
                "reason": f"описание содержит ключевые слова: {', '.join([k for k in keywords if k in video_note])}"
            })
            continue
        
        # Проверяем общее соответствие по вопросу
        if any(word in video_title or word in video_note 
               for word in question_lower.split() if len(word) > 3):
            suitable_videos.append({
                "title": video["title"],
                "reason": "содержит релевантную информацию по теме вопроса"
            })
    
    return {"suitable_videos": suitable_videos}

def extract_keywords_from_question(question: str) -> list:
    """Извлекает ключевые слова из вопроса для поиска релевантных материалов"""
    question_lower = question.lower()
    
    # Список технических терминов, которые могут быть в вопросах
    technical_terms = [
        "bias", "variance", "overfitting", "underfitting", "регуляризация", "валидация", "dropout",
        "precision", "recall", "f1", "roc", "auc", "pr", "дисбаланс", "классов",
        "k-fold", "cross-validation", "разбиение", "среднее",
        "leakage", "стратификация", "shuffle", "random_state",
        "scaling", "encoding", "interaction", "feature", "engineering",
        "машинное обучение", "ml", "алгоритм", "модель", "обучение",
        "нейронная сеть", "дерево решений", "классификация", "регрессия",
        "python", "последовательности", "sequences", "zen", "философия",
        "задание", "домашнее", "дз", "контрольная", "тест", "лекция"
    ]
    
    found_keywords = []
    for term in technical_terms:
        if term in question_lower:
            found_keywords.append(term)
    
    # Также добавляем слова из вопроса, если они длиннее 3 символов
    question_words = [word for word in question_lower.split() if len(word) > 3]
    found_keywords.extend(question_words)
    
    # Убираем дубликаты
    return list(set(found_keywords))

def find_matching_assignments(question: str, assignments: list) -> list:
    """Находит задания, которые соответствуют вопросу"""
    question_lower = question.lower()
    matching = []
    
    for assignment in assignments:
        # Проверяем соответствие по названию задания
        if any(word in assignment["title"].lower() for word in question_lower.split()):
            matching.append(assignment)
            continue
        
        # Проверяем соответствие по вопросам в задании
        for q in assignment.get("questions", []):
            if any(word in q["title"].lower() for word in question_lower.split()):
                matching.append(assignment)
                break
            
            # Проверяем ключевые слова
            for keyword in q.get("keywords", []):
                if isinstance(keyword, dict):
                    keyword_text = keyword.get("word", "").lower()
                else:
                    keyword_text = str(keyword).lower()
                
                if keyword_text in question_lower:
                    matching.append(assignment)
                    break
    
    return matching

def handle_overdue_question(question: str, subject: str, context: dict) -> str:
    """Обрабатывает вопросы о просроченных заданиях"""
    today = datetime.now()
    overdue_assignments = []
    
    for assignment in context["assignments"]:
        due_date_str = assignment.get("due_date", "")
        if due_date_str:
            try:
                due_date = datetime.fromisoformat(due_date_str)
                days_left = (due_date - today).days
                
                if days_left < 0:
                    overdue_assignments.append({
                        "title": assignment["title"],
                        "days_overdue": abs(days_left),
                        "questions": assignment["questions"]
                    })
            except:
                pass
    
    if not overdue_assignments:
        return "✅ У вас нет просроченных заданий!"
    
    response_parts = ["🔴 **Просроченные задания:**"]
    
    for assignment in overdue_assignments:
        response_parts.append(f"\n**{assignment['title']}** (просрочено на {assignment['days_overdue']} дн.)")
        
        # Рекомендуем материалы для каждого вопроса
        for question in assignment["questions"]:
            keywords = []
            for kw in question.get("keywords", []):
                if isinstance(kw, dict):
                    keywords.append(kw.get("word", ""))
                else:
                    keywords.append(str(kw))
            
            if keywords:
                response_parts.append(f"  • {question['title']} - ключевые слова: {', '.join(keywords)}")
    
    response_parts.append("\n💡 **Рекомендации:**")
    response_parts.append("• Свяжитесь с преподавателем для уточнения возможности сдачи")
    response_parts.append("• Изучите материалы курса для подготовки ответов")
    response_parts.append("• Обратитесь к чат-ассистенту за помощью с конкретными вопросами")
    
    return "\n".join(response_parts)

def handle_general_question(question: str, subject: str, context: dict) -> str:
    """Обрабатывает общие вопросы через AI"""
    if client is None:
        return "❌ AI недоступен для ответа на общие вопросы"
    
    try:
        # Формируем контекст для AI
        context_text = f"""
Предмет: {subject}
Доступные материалы: {len(context['materials'])} файлов
Доступные видео: {len(context['videos'])} видео
Активных заданий: {len(context['assignments'])}
"""
        
        system_msg = (
            "Ты - помощник студента по учебному курсу. "
            "Отвечай на вопросы студента, помогай с навигацией по курсу. "
            "Если вопрос касается конкретных заданий или материалов, рекомендуй соответствующие ресурсы. "
            "Отвечай на русском языке, будь полезным и конкретным."
        )
        
        user_msg = f"""
Контекст курса:
{context_text}

Вопрос студента: {question}

Помоги студенту с его вопросом.
"""
        
        resp = client.chat.completions.create(
            model=DEEPSEEK_MODEL,
            temperature=0.4,
            messages=[
                {"role": "system", "content": system_msg},
                {"role": "user", "content": user_msg}
            ],
        )
        
        return resp.choices[0].message.content.strip()
        
    except Exception as e:
        return f"❌ Ошибка обработки вопроса: {str(e)}"



# Инициализация JWT токена
if "jwt_token" not in st.session_state:
    st.session_state["jwt_token"] = None

# Проверка авторизации - если пользователь не авторизован, показываем страницу входа
require_auth()

st.set_page_config(page_title="EduAI Hub (Mini)", page_icon="🎓", layout="wide")
st.markdown("""
<style>
/* Layout spacing */
.block-container { padding-top: 2rem !important; padding-bottom: 2rem !important; }

/* Cards */
.card {
  border: 1px solid #eaeaea;
  border-radius: 16px;
  padding: 16px 16px;
  box-shadow: 0 2px 8px rgba(0,0,0,0.04);
  background: #fff;
  margin-bottom: 12px;
}

/* Buttons */
div.stButton > button {
  border-radius: 10px;
  padding: 0.55rem 1rem;
  border: 1px solid #e5e7eb;
}

/* Badges (chips) */
.badge {
  display: inline-block;
  padding: 2px 8px;
  border-radius: 999px;
  background: #eef2ff;
  color: #3730a3;
  font-size: 12px;
  border: 1px solid #e0e7ff;
}

/* Muted small text */
.muted { color: #6b7280; font-size: 12px; }

/* Tables (slightly denser) */
.dataframe tbody tr th, .dataframe tbody tr td { padding: 6px 10px; }
</style>
""", unsafe_allow_html=True)

st.title("EduAI Hub — платформа для эффективного обучения")
st.divider()


# --- tiny "db" with persistence
if "points" not in st.session_state:
    st.session_state.points = load_json(POINTS_PATH, {})
if "submissions" not in st.session_state:
    st.session_state.submissions = load_json(SUBM_PATH, [])
if "reviews" not in st.session_state:
    st.session_state.reviews = load_json(REVIEWS_PATH, [])

migrate_stores_to_subject_scope()

# --- load rubric
with open("rubric.json", "r", encoding="utf-8") as f:
    RUBRICS = json.load(f)   # list of rubrics (>=1)
id_to_rubric = {r["assignment_id"]: r for r in RUBRICS}

def keyword_score_for(text: str, kw_list: list[dict], max_points: int):
    text_lc = (text or "").lower()
    score = 0
    details = []
    for kw in kw_list:
        found = bool(re.search(rf"\b{re.escape(kw['word'].lower())}\b", text_lc))
        if found:
            score += kw["points"]
            details.append(f"+{kw['points']}: найдено '{kw['word']}'")
        else:
            details.append(f"+0: нет '{kw['word']}'")
    score = min(score, max_points)
    return score, details


def llm_grade(answer_text: str, rubric: dict) -> dict | None:
    """
    Ask DeepSeek to grade 0–100 and give short feedback bullets.
    Returns {"score": int, "feedback": [str, ...]} or None on failure.
    """
    if client is None:
        st.warning("AI-оценка: не найден ключ (DEEPSEEK_API_KEY).")
        return None
    try:
        system_msg = (
            "You are a strict but fair TA. "
            "Grade the student's short answer on a 0–100 scale based ONLY on the rubric. "
            "Return strict JSON: {\"score\": <int>, \"feedback\": [\"...\", \"...\"]}"
        )
        user_msg = (
            "Rubric title: " + rubric.get("title","") + "\n"
            "Max points: 100\n"
            "Keywords (hints): " + ", ".join([k['word'] for k in rubric.get('keywords', [])]) + "\n"
            "Student answer:\n" + answer_text
        )

        resp = client.chat.completions.create(
            model=DEEPSEEK_MODEL,
            temperature=0.2,
            messages=[
                {"role": "system", "content": system_msg},
                {"role": "user", "content": user_msg}
            ],
        )
        content = resp.choices[0].message.content.strip()

        # Extract JSON
        import json, re
        m = re.search(r"\{.*\}", content, re.S)
        json_str = m.group(0) if m else content
        data = json.loads(json_str)

        score = int(max(0, min(100, data.get("score", 0))))
        feedback = data.get("feedback", [])
        if not isinstance(feedback, list): feedback = [str(feedback)]
        return {"score": score, "feedback": feedback[:5]}

    except Exception as e:
        # Surface the real reason in the UI so you know what to fix
        error_msg = f"AI ошибка: {getattr(e, 'message', str(e))[:200]}"
        st.error(error_msg)
        return None

def award_points(user: str, pts: int):
    pts_store = get_points_store()
    pts_store[user] = pts_store.get(user, 0) + pts

def peer_avg_for_submission(subm_idx: int):
    reviews = get_reviews_store()
    scores = [r["avg_score"] for r in reviews if r["submission_idx"] == subm_idx]
    return (round(mean(scores), 2), len(scores)) if scores else (None, 0)

# --- sidebar "login"
st.sidebar.markdown("### ⚙️ Настройки")
st.sidebar.header("Профиль")

# Получаем данные текущего пользователя
current_username = get_current_username()
current_role = get_current_role()
role_display = "Студент" if current_role == "student" else "Преподаватель"

st.sidebar.markdown(f"**👤 {current_username}**")
if current_role == "student":
    st.sidebar.markdown(f"**🎓 {role_display}**")
else:
    st.sidebar.markdown(f"**👨‍🏫 {role_display}**")


if st.sidebar.button("🚪 Выйти", use_container_width=True):
    logout_user()
    st.rerun()

st.sidebar.divider()


user = current_username

# --- subject picker
# Загружаем предметы из файла subjects.json
def load_subjects():
    try:
        with open("subjects.json", "r", encoding="utf-8") as f:
            subjects = json.load(f)
        return subjects
    except FileNotFoundError:
        # Если файл не существует, создаем его с предметами из rubric.json
        subjects = sorted({ r.get("subject", "General") for r in RUBRICS })
        save_subjects(subjects)
        return subjects

def save_subjects(subjects):
    try:
        with open("subjects.json", "w", encoding="utf-8") as f:
            json.dump(subjects, f, ensure_ascii=False, indent=2)
    except Exception as e:
        st.error(f"Ошибка сохранения предметов: {e}")

SUBJECTS = load_subjects()

# Функционал управления предметами (только для преподавателей)
if current_role == "teacher":
    st.sidebar.markdown("### 📚 Управление предметами")
    
    # Поле для ввода нового предмета
    new_subject = st.sidebar.text_input("Новый предмет", placeholder="Введите название предмета", key="new_subject_input")
    
    
    if st.sidebar.button("➕ Добавить предмет", use_container_width=True):
        if new_subject and new_subject.strip():
            new_subject = new_subject.strip()
            if new_subject not in SUBJECTS:
                SUBJECTS.append(new_subject)
                SUBJECTS.sort()
                save_subjects(SUBJECTS)
                st.session_state.subject = new_subject
                st.sidebar.success(f"✅ Предмет '{new_subject}' добавлен!")
                st.rerun()
            else:
                st.sidebar.warning("⚠️ Предмет уже существует")
        else:
            st.sidebar.error("❌ Введите название предмета")
    
    if len(SUBJECTS) > 1:  
        st.sidebar.markdown("**Удалить предмет:**")
        
        # Создаем список предметов для удаления (исключаем General)
        subjects_to_delete = [s for s in SUBJECTS if s != "General"]
        
        if subjects_to_delete:
            subject_to_delete = st.sidebar.selectbox(
                "Выберите предмет для удаления",
                subjects_to_delete,
                key="subject_to_delete",
                help="Предмет можно удалить только если в нем нет заданий"
            )
            
            if st.sidebar.button("🗑️ Удалить предмет", use_container_width=True):
                assignments_in_subject = [r for r in RUBRICS if r.get("subject") == subject_to_delete]
                if not assignments_in_subject:
                    SUBJECTS.remove(subject_to_delete)
                    save_subjects(SUBJECTS)
                    if st.session_state.get("subject") == subject_to_delete:
                        st.session_state.subject = "General"
                    st.sidebar.success(f"✅ Предмет '{subject_to_delete}' удален!")
                    st.rerun()
                else:
                    st.sidebar.error(f"❌ Нельзя удалить предмет с заданиями ({len(assignments_in_subject)} заданий)")
    
    st.sidebar.markdown("---")

subject = st.sidebar.selectbox("Предмет", SUBJECTS, key="subject")

use_ai = st.sidebar.checkbox("Включить AI-оценку (опционально)", value=True)
st.session_state["use_ai"] = bool(use_ai)  

# Connection status
if client is None:
    st.sidebar.error("AI отключен: нет DEEPSEEK_API_KEY в .env")
else:
    try:
        _ = client.models.list()  
        st.sidebar.success("AI подключен (DeepSeek)")
    except Exception as e:
        st.sidebar.error(f"AI недоступен: {getattr(e, 'message', str(e))[:120]}")
st.sidebar.info("Добро пожаловать в EduAI Hub!")


st.sidebar.divider()
st.sidebar.markdown("**О проекте**\n\nОбразовательная платформа с AI-ассистентом для создания тестов, проверки заданий, кросс-оценки и управления материалами курса.")


if current_role == "teacher":
    tab_submit, tab_materials, tab_leaderboard, tab_peer, tab_chat, tab_test_creator = st.tabs(
        ["📝 Мои задания", "📚 Материалы курса", "🏆 Лидерборд", "🤝 Кросс-проверка", "💬 Чат-ассистент (демо)", "📝 Создание тестов"]
    )
else:
    tab_submit, tab_materials, tab_leaderboard, tab_peer, tab_chat = st.tabs(
        ["📝 Мои задания", "📚 Материалы курса", "🏆 Лидерборд", "🤝 Кросс-проверка", "💬 Чат-ассистент (демо)"]
    )

with tab_peer:

    st.subheader("Кросс-проверка (анонимно)")
    st.subheader("Выбор задания")
    RUBRIC = assignment_selector("assign_id_peer")
    if RUBRIC is None:
        st.info("Добавьте задание в этом предмете во вкладке 'Создание тестов'.")
    else:
        st.subheader("Ответы на вопросы задания")

        # Pick which assignment to review (reuse current selection)
        st.caption("Выберите задание для проверки (совпадает с селектором выше).")
        # We already have assign_id and RUBRIC bound from the selector section.

        # Build candidate submissions
        subs = get_submissions_store()
        all_subms = [
            (idx, s) for idx, s in enumerate(subs)
            if s.get("assignment") == RUBRIC["assignment_id"] and s.get("user") != user
        ]

        if not all_subms:
            st.info("Пока нет чужих ответов по этому заданию.")
        else:
            # Show a simple picker
            option_indices = [idx for idx, _ in all_subms]
            labels = [f"#{idx} — {s.get('user','?')}: {submission_preview(s)}" for idx, s in all_subms]
            sel = st.selectbox("Выберите ответ для проверки:", options=range(len(option_indices)),
                            format_func=lambda i: labels[i])
            pick_idx = option_indices[sel]
            pick = subs[pick_idx]


            st.write("Ответ студента")
            if "answer" in pick:
                st.write(pick["answer"])
            elif "answers" in pick and pick["answers"]:
                for ans in pick["answers"]:
                    qid = ans.get("question_id", "")
                    atxt = (ans.get("answer") or "").strip()
                    st.markdown(f"**{qid}**")
                    st.write(atxt if atxt else "—")
                    st.markdown("<hr style='border:none;border-top:1px solid #eee;margin:8px 0;' />", unsafe_allow_html=True)
            else:
                st.write("Ответ отсутствует.")


            st.markdown("### Оцените по критериям (1–5)")
            c1, c2 = st.columns(2)
            with c1:
                sc_relevance  = st.slider("Соответствие заданию", 1, 5, 3)
                sc_structure  = st.slider("Структура и логика", 1, 5, 3)
            with c2:
                sc_argument   = st.slider("Аргументация / примеры", 1, 5, 3)
                sc_clarity    = st.slider("Ясность изложения", 1, 5, 3)
            comment = st.text_area("Короткий комментарий (опционально)", placeholder="Что можно улучшить?")

            avg_score = round((sc_relevance + sc_structure + sc_argument + sc_clarity) / 4, 2)
            st.info(f"Средняя оценка по критериям: **{avg_score} / 5**")

            if RUBRIC and st.button("Отправить отзыв"):
                review = {
                    "submission_idx": pick_idx,
                    "assignment": RUBRIC["assignment_id"],
                    "reviewer": user,          
                    "scores": {
                        "relevance": sc_relevance,
                        "structure": sc_structure,
                        "argument": sc_argument,
                        "clarity": sc_clarity
                    },
                    "avg_score": avg_score,
                    "comment": comment.strip()
                }
                reviews = get_reviews_store()
                reviews.append(review)
                save_reviews_store()

                award_points(user, 1)
                save_points_store()
                st.success("Отзыв сохранен! Вам начислено +1 очко за кросс-проверку.")

    st.markdown("---")
    st.subheader("Мои полученные отзывы")
    subs = get_submissions_store()
    my_subms = [] if RUBRIC is None else [
        (i, s) for i, s in enumerate(subs)
        if s.get("assignment") == RUBRIC["assignment_id"] and s.get("user") == user
    ]
    reviews = get_reviews_store()

    if not my_subms:
        st.write("Вы ещё не сдавали это задание.")
    else:
        for idx, s in my_subms:
            avg, cnt = peer_avg_for_submission(idx)
            st.markdown(
                f"**Сдача #{idx}** — peer-оценок: {cnt}"
                + (f", среднее: **{avg}/5**" if avg is not None else "")
            )
            comments = [
                r.get("comment") for r in reviews
                if isinstance(r, dict)
                and r.get("submission_idx") == idx
                and r.get("comment")
            ]
            if comments:
                with st.expander("Комментарии"):
                    for c in comments[-5:]:
                        st.write("•", c)


with tab_submit:

    st.subheader("Выбор задания")
    RUBRIC = assignment_selector("assign_id_submit")

    if RUBRIC is None:
        st.info("В этом предмете ещё нет заданий. Создайте тест во вкладке 'Создание тестов'.")
    else:
        # RUBRIC is bound from the assignment selector (as before)
        questions = RUBRIC.get("questions", [])
        # Информация о тесте и кнопка старта
        due = RUBRIC.get("due_date")
        tlim = RUBRIC.get("time_limit_minutes")
        info_cols = st.columns(3)
        info_cols[0].markdown(f"**Вопросов:** {len(questions)}")
        info_cols[1].markdown(f"**Дедлайн:** {due if due else '—'}")
        time_box = info_cols[2].empty()
        time_box.markdown(f"**Время:** {tlim} мин" if tlim else "**Время:** не ограничено")

        start_key = f"start_{RUBRIC['assignment_id']}"
        started = st.session_state.get(start_key, False)
        start_ts_key = f"start_ts_{RUBRIC['assignment_id']}"
        time_up_key = f"time_up_{RUBRIC['assignment_id']}"
        if not started:
            if st.button("Начать", key=f"btn_start_{RUBRIC['assignment_id']}"):
                st.session_state[start_key] = True
                st.session_state[start_ts_key] = time.time()
                st.session_state[time_up_key] = False
                st.rerun()
        else:
            # Таймер, если есть лимит (обновляем, но не ререндрим до конца блока рендера)
            if tlim and int(tlim) > 0:
                if start_ts_key not in st.session_state:
                    st.session_state[start_ts_key] = time.time()
                elapsed = int(time.time() - st.session_state[start_ts_key])
                total = int(tlim) * 60
                remaining = max(0, total - elapsed)
                mm = remaining // 60
                ss = remaining % 60
                time_box.markdown(f"**Осталось:** {mm:02d}:{ss:02d}")
                if remaining == 0:
                    st.session_state[time_up_key] = True
            if not questions:
                st.info("Для этого задания ещё не настроены вопросы.")
            else:
                test_type = RUBRIC.get("test_type")
                if not test_type:
                    if any("options" in q for q in questions):
                        test_type = "multiple_choice"
                    else:
                        test_type = "keyword_based"
                
                time_up = st.session_state.get(time_up_key, False)
                test_completed = st.session_state.get(f"completed_{RUBRIC['assignment_id']}", False)
                
                # Если тест завершен, показываем результаты
                if test_completed:
                    st.success("✅ Тест завершен!")
                    
                    # Получаем последний результат из submissions
                    subs = get_submissions_store()
                    last_submission = None
                    for sub in reversed(subs):
                        if (sub.get("user") == user and 
                            sub.get("assignment") == RUBRIC["assignment_id"]):
                            last_submission = sub
                            break
                    
                    if last_submission:
                        st.markdown(f"**Итоговый результат: {last_submission.get('total_score', 0)}/{last_submission.get('total_max', 0)}**")
                        st.info(f"Начислено очков: {last_submission.get('points_awarded', 0)}")
                        
                        # Показываем детали по вопросам
                        for answer in last_submission.get("answers", []):
                            q_id = answer.get("question_id")
                            q_title = next((q["title"] for q in questions if q["question_id"] == q_id), q_id)
                            score = answer.get("score", 0)
                            
                            with st.expander(f"{q_id}: {q_title} — {score} баллов"):
                                st.write(f"**Ваш ответ:** {answer.get('answer', '')}")
                                st.write(f"**Итоговые баллы:** {score}")
                                
                                # Показываем детали оценки
                                ai_score = answer.get("ai_score")
                                if ai_score is not None:
                                    # Вычисляем оценку по ключевым словам
                                    question = next((q for q in questions if q["question_id"] == q_id), None)
                                    if question:
                                        kw_score = score - (ai_score * 0.3 * question["max_points"] / 100) / 0.7
                                        kw_score = round(kw_score)
                                        st.write(f"**Баллы по ключевым словам:** {kw_score}")
                                        st.write(f"**AI оценка:** {ai_score}/100")
                                        st.write(f"**Формула:** {kw_score} × 0.7 + {ai_score} × 0.3 = {score}")
                                
                                # Показываем AI фидбек, если есть
                                ai_feedback = answer.get("ai_feedback", [])
                                if ai_feedback:
                                    st.markdown("**🤖 AI-фидбек:**")
                                    for fb in ai_feedback:
                                        st.write(f"• {fb}")
                        
                        st.markdown("---")
                        if st.button("🔄 Начать тест заново", type="secondary"):
                            # Очищаем состояние для нового теста
                            st.session_state[start_key] = False
                            st.session_state[time_up_key] = False
                            st.session_state[f"completed_{RUBRIC['assignment_id']}"] = False
                            if start_ts_key in st.session_state:
                                del st.session_state[start_ts_key]
                            # Очищаем ответы
                            for q in questions:
                                answer_key = f"ans_{RUBRIC['assignment_id']}_{q['question_id']}"
                                if answer_key in st.session_state:
                                    del st.session_state[answer_key]
                            st.rerun()
                    else:
                        st.error("Результаты теста не найдены")
                        if st.button("🔄 Начать тест заново", type="secondary"):
                            st.session_state[f"completed_{RUBRIC['assignment_id']}"] = False
                            st.rerun()
                else:
                    # Формируем ответы (даже если время вышло, чтобы можно было автозавершить)
                    answers = {}
                    
                    if test_type == "multiple_choice":
                        # Render multiple choice questions
                        for q in questions:
                            st.markdown(f"### {q['title']}  \n*Макс. баллов: {q['max_points']}*")
                            
                            options = q.get("options", [])
                            if options:
                                selected_option = st.radio(
                                    "Выберите правильный ответ:",
                                    options=options,
                                    index=None,  # Никакой вариант не выбран по умолчанию
                                    key=f"ans_{RUBRIC['assignment_id']}_{q['question_id']}"
                                )
                                answers[q["question_id"]] = selected_option
                            else:
                                st.warning("Нет вариантов ответов для этого вопроса")
                                answers[q["question_id"]] = ""
                            
                            st.markdown("---")
                    else:
                        # Render text area for keyword-based questions (original logic)
                        for q in questions:
                            st.markdown(f"### {q['title']}  \n*Макс. баллов: {q['max_points']}*")
                            answers[q["question_id"]] = st.text_area(
                                f"Ваш ответ ({q['question_id']})",
                                key=f"ans_{RUBRIC['assignment_id']}_{q['question_id']}",
                                height=140
                            )
                            st.markdown("---")

                    # Завершение: либо по кнопке, либо автоматически по таймеру
                    do_finish = False
                    if time_up:
                        st.warning("Время выполнения теста истекло. Ответы отправлены автоматически.")
                        # Подхватим значения из session_state, если локально пусто
                        for q in questions:
                            key = f"ans_{RUBRIC['assignment_id']}_{q['question_id']}"
                            if q["question_id"] not in answers or answers[q["question_id"]] == "":
                                answers[q["question_id"]] = st.session_state.get(key, "")
                        do_finish = True
                    else:
                        # Проверяем, что все вопросы отвечены
                        all_answered = True
                        for q in questions:
                            if test_type == "multiple_choice":
                                if answers.get(q["question_id"]) is None:
                                    all_answered = False
                                    break
                            else:
                                if not answers.get(q["question_id"], "").strip():
                                    all_answered = False
                                    break
                        
                        if all_answered:
                            do_finish = st.button("Завершить", type="primary")
                        else:
                            st.warning("⚠️ Ответьте на все вопросы перед завершением")
                            do_finish = False

                    if do_finish:
                        total_score = 0
                        total_max = sum(q["max_points"] for q in questions)
                        per_q_results = []

                        # Run scoring per question based on test type
                        for q in questions:
                            ans_text = answers.get(q["question_id"], "")
                            
                            if test_type == "multiple_choice":
                                correct_answer = q.get("correct_answer", "")
                                if ans_text == correct_answer:
                                    q_score = q["max_points"]
                                    q_details = [f"✅ Правильный ответ: +{q['max_points']} баллов"]
                                else:
                                    q_score = 0
                                    q_details = [f"❌ Неправильный ответ: 0 баллов (правильный: {correct_answer})"]
                                
                                per_q_results.append({
                                    "question_id": q["question_id"],
                                    "title": q["title"],
                                    "answer": ans_text,
                                    "kw_score": q_score,
                                    "details": q_details
                                })
                                total_score += q_score
                                
                            else:
                                q_score_kw, q_details = keyword_score_for(
                                    ans_text, q["keywords"], q["max_points"]
                                )
                                per_q_results.append({
                                    "question_id": q["question_id"],
                                    "title": q["title"],
                                    "answer": ans_text,
                                    "kw_score": q_score_kw,
                                    "details": q_details
                                })
                                total_score += q_score_kw

                        # Optional AI refinement (only for keyword-based tests)
                        if test_type == "keyword_based":
                            use_ai = st.session_state.get("use_ai")
                            if 'use_ai' not in st.session_state:
                                try:
                                    use_ai_local = use_ai
                                except:
                                    use_ai_local = False
                            else:
                                use_ai_local = st.session_state['use_ai']

                            ai_total_delta = 0
                            if use_ai_local:
                                for item in per_q_results:
                                    if client is None:
                                        continue
                                    # Build a temporary one-question rubric for llm_grade
                                    one_q_rubric = {
                                        "title": f"{RUBRIC['title']} — {item['title']}",
                                        "keywords": [{"word": k["word"]} for k in next(q for q in questions if q["question_id"]==item["question_id"])["keywords"]]
                                    }
                                    with st.spinner(f"AI оценивает: {item['question_id']}"):
                                        llm_info = llm_grade(item["answer"], one_q_rubric)
                                    if llm_info:
                                        # Combine keyword score with AI score using weighted average
                                        # 70% weight for keywords, 30% weight for AI
                                        question_max = next(q for q in questions if q["question_id"]==item["question_id"])["max_points"]
                                        kw_pct = (item["kw_score"] / question_max) * 100
                                        ai_pct = llm_info["score"]
                                        
                                        # Weighted average: 70% keywords + 30% AI
                                        combined_pct = (kw_pct * 0.7) + (ai_pct * 0.3)
                                        combined_score = round(combined_pct / 100 * question_max)
                                        
                                        ai_total_delta += (combined_score - item["kw_score"])
                                        item["ai_score"] = int(round(llm_info["score"]))
                                        item["final_score"] = combined_score
                                        item["ai_feedback"] = llm_info.get("feedback", [])
                                    else:
                                        item["final_score"] = item["kw_score"]
                                total_score += ai_total_delta  # adjust total if AI used
                            else:
                                for item in per_q_results:
                                    item["final_score"] = item["kw_score"]
                        else:
                            for item in per_q_results:
                                item["final_score"] = item["kw_score"]

                        # Show per-question breakdown
                        st.success(f"Итог: {total_score}/{total_max}")
                        for item in per_q_results:
                            final_sc = item.get("final_score", item.get("kw_score", 0))
                            kw_sc = item.get("kw_score", 0)
                            ai_sc = item.get("ai_score")
                            
                            with st.expander(f"{item['question_id']}: {item['title']} — {final_sc} баллов"):
                                st.markdown("**По ключевым словам:**")
                                for d in item["details"]:
                                    st.write("•", d)
                                st.write(f"**Баллы по ключевым словам:** {kw_sc}")
                                
                                if ai_sc is not None:
                                    st.write(f"**AI оценка:** {ai_sc}/100")
                                    st.write(f"**Итоговая оценка:** {kw_sc} × 0.7 + {ai_sc} × 0.3 = {final_sc}")
                                
                                if item.get("ai_feedback"):
                                    st.markdown("**AI-фидбек:**")
                                    for fb in item["ai_feedback"]:
                                        st.write("•", fb)

                        # Award points (same rule: 1 балл за каждые ~10% итога, минимум 1)
                        pct = total_score / max(1, total_max)
                        gained = max(1, int(round(pct * 10)))
                        award_points(user, gained)
                        save_points_store()

                        # Save a multi-question submission
                        subs = get_submissions_store()
                        subs.append({
                            "user": user,
                            "assignment": RUBRIC["assignment_id"],
                            "answers": [
                                {
                                    "question_id": item["question_id"],
                                    "answer": item["answer"],
                                    "score": item.get("final_score", item.get("kw_score", 0)),
                                    "ai_feedback": item.get("ai_feedback", []),
                                    "ai_score": item.get("ai_score", None)
                                } for item in per_q_results
                            ],
                            "total_score": total_score,
                            "total_max": total_max,
                            "points_awarded": gained
                        })
                        save_submissions_store()

                        st.info(f"Начислено {gained} очк.")
                        
                        # Помечаем тест как завершенный
                        st.session_state[f"completed_{RUBRIC['assignment_id']}"] = True
                        st.session_state[time_up_key] = True
                        
                        st.success("✅ Тест завершен! Результаты сохранены.")
                        st.rerun()

            # В самом конце: если есть таймер и он не истёк — плавно ререндерим раз в секунду
            if tlim and int(tlim) > 0 and not st.session_state.get(time_up_key, False):
                time.sleep(1)
                st.rerun()


with tab_leaderboard:
    st.subheader(f"Лидерборд — {current_subject()}")
    pts = get_points_store()
    if not pts:
        st.write("Пока пусто.")
    else:
        board = sorted(pts.items(), key=lambda x: x[1], reverse=True)
        for i, (u, p) in enumerate(board, start=1):
            st.write(f"{i}. **{u}** — {p} очк. {badge_for(p)}")

     # описание уровней
    st.markdown("### 🏅 Значки и уровни")
    st.markdown("""
    - 🌱 **Rookie** — 0–4 очка. Начало пути!
    - 🥉 **Starter** — 5–14 очков. Первые шаги.
    - 🥈 **Getting There** — 15–29 очков. Хороший прогресс.
    - 🥇 **Consistent** — 30–49 очков. Стабильные результаты.
    - 🏅 **Power Learner** — 50+ очков. Отличная вовлечённость!
    """)

    subs = get_submissions_store()
    with st.expander("Последние сдачи"):
        if not subs:
            st.write("Нет сдач.")
        else:
            for s in reversed(subs[-10:]):
                summary = f"{s['total_score']}/{s['total_max']}" if "total_score" in s else str(s.get("score","?"))
                preview = submission_preview(s, max_len=40)
                st.markdown(f"- **{s['user']}** → {summary}  <span class='muted'>(+{s['points_awarded']} очк.)</span>  — {preview}", unsafe_allow_html=True)

    pts = get_points_store()
    subs = get_submissions_store()
    revs = get_reviews_store()

    if pts:
        df_points = pd.DataFrame([{"user": u, "points": p} for u, p in pts.items()]).sort_values(
            "points", ascending=False
        )
    else:
        df_points = pd.DataFrame(columns=["user", "points"])

    df_subm = pd.DataFrame(subs) if subs else pd.DataFrame(columns=[
        "user", "assignment", "answers", "total_score", "total_max", "points_awarded"
    ])

    df_rev = pd.DataFrame(revs) if revs else pd.DataFrame(columns=[
        "submission_idx", "assignment", "reviewer", "scores", "avg_score", "comment"
    ])

    st.write("—")
    st.markdown("**Выгрузка данных**")
    buf1, buf2, buf3 = io.BytesIO(), io.BytesIO(), io.BytesIO()
    df_points.to_csv(buf1, index=False, encoding="utf-8-sig")
    df_subm.to_csv(buf2, index=False, encoding="utf-8-sig")
    df_rev.to_csv(buf3, index=False, encoding="utf-8-sig")
    st.download_button("⬇️ Лидерборд (CSV)", data=buf1.getvalue(),
                    file_name=f"leaderboard_{current_subject()}.csv", mime="text/csv")
    st.download_button("⬇️ Сдачи (CSV)", data=buf2.getvalue(),
                    file_name=f"submissions_{current_subject()}.csv", mime="text/csv")
    st.download_button("⬇️ Отзывы (CSV)", data=buf3.getvalue(),
                    file_name=f"reviews_{current_subject()}.csv", mime="text/csv")


with tab_chat:
    st.subheader("🤖 Чат-ассистент")
    st.caption("AI-помощник для навигации по курсу, проверки дедлайнов и рекомендаций материалов.")
    
    # Раскрывающаяся панель с краткой информацией
    with st.expander("📋 Краткая информация по курсу", expanded=False):
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("**📝 Активные задания:**")
            current_subj = current_subject()
            subj_rubrics = [r for r in RUBRICS if r.get("subject", "General") == current_subj]
            
            if subj_rubrics:
                for rubric in subj_rubrics:  # Показываем все задания
                    due = rubric.get("due_date")
                    if due:
                        try:
                            from datetime import datetime
                            due_date = datetime.fromisoformat(due)
                            today = datetime.now()
                            days_left = (due_date - today).days
                            
                            if days_left < 0:
                                status = "🔴 Просрочено"
                            elif days_left == 0:
                                status = "🟡 Сегодня"
                            elif days_left <= 3:
                                status = f"🟠 Осталось {days_left} дн."
                            else:
                                status = f"🟢 Осталось {days_left} дн."
                        except:
                            status = f"📅 до {due}"
                    else:
                        status = "⏰ Без дедлайна"
                    
                    st.markdown(f"• **{rubric['title']}** - {status}")
            else:
                st.markdown("*Нет активных заданий*")
        
        with col2:
            st.markdown("**📚 Доступные материалы:**")
            idx, materials = get_materials_store()
            if materials:
                for material in materials[:5]:  # Показываем только первые 5
                    display_name = display_filename(material['name'], material.get('original_name'))
                    st.markdown(f"• **{display_name}**")
                    if material.get("note"):
                        st.markdown(f"  *{material['note']}*")
            else:
                st.markdown("*Нет материалов*")
    
    # Блок с примерными вопросами
    st.markdown("### 💡 Примерные вопросы:")
    
    example_questions = [
        "Какие задания у меня просрочены?",
        "Какие материалы есть по машинному обучению?",
        "Когда дедлайн по заданию 2?"
    ]
    
    # Отображаем примерные вопросы как простой текст
    for question in example_questions:
        st.markdown(f"• {question}")
    
    # Поле для ввода вопроса
    st.markdown("### 💬 Задайте свой вопрос:")
    
    user_question = st.text_input(
        "Введите ваш вопрос:",
        placeholder="Например: Какие материалы помогут с заданием по ML?",
        key="user_question_input"
    )
    
    # Кнопка отправки
    if st.button("🚀 Отправить вопрос", use_container_width=True, type="primary"):
        if user_question.strip():
            # Здесь будет логика обработки вопроса через нейронную сеть
            st.session_state.current_question = user_question.strip()
            st.rerun()
        else:
            st.warning("Пожалуйста, введите вопрос")
    
    # Область для ответа
    if hasattr(st.session_state, 'current_question') and st.session_state.current_question:
        st.markdown("### 🤖 Ответ ассистента:")
        
        # Показываем вопрос пользователя
        st.markdown(f"**Ваш вопрос:** {st.session_state.current_question}")
        st.markdown("---")
        
        with st.spinner("Анализирую ваш вопрос..."):
            # Здесь будет вызов функции обработки вопроса
            response = process_chat_question(st.session_state.current_question, current_subj)
        
        # Отображаем ответ в красивом контейнере
        with st.container(border=True):
            st.markdown(response)
        
        # Показываем дополнительную информацию
        st.markdown("---")
        st.markdown("💡 **Совет:** Вы можете задавать вопросы о дедлайнах, материалах курса, заданиях и получать персонализированные рекомендации.")

st.caption("🚀 Полнофункциональная платформа: AI-генерация тестов, умная оценка, кросс-проверка, чат-ассистент, управление материалами.")

with tab_materials:
    st.subheader(f"Материалы: {current_subject()}")
    tab_documents, tab_videos = st.tabs(["📄 Печатные материалы", "🎥 Видео материалы"])
    
    with tab_documents:
        st.markdown("### 📄 Печатные материалы")
        
        # Upload area 
        if current_role == "teacher":
            st.caption("Загрузите один или несколько файлов. Они будут доступны в этом предмете.")
            uploaded = st.file_uploader("Загрузить файлы", type=None, accept_multiple_files=True)
            note = st.text_input("Короткое описание (опционально)", placeholder="Например: лекция 1, слайды")
            col_u1, col_u2 = st.columns([1, 2])
            do_save = col_u1.button("Загрузить")
        else:
            st.info("📚 Загрузка материалов доступна только преподавателям")
            uploaded = None
            note = ""
            do_save = False

        # Логика загрузки и отображения печатных материалов
        idx, items = get_materials_store()
        subj = current_subject()
        subj_dir = os.path.join(MATERIALS_DIR, safe_filename(subj))
        ensure_dir(subj_dir)

        if do_save and uploaded:
            import time
            for f in uploaded:
                raw = f.read()
                fname = safe_filename(f.name)
                path = os.path.join(subj_dir, fname)

                # avoid accidental overwrite: add suffix if exists
                base, ext = os.path.splitext(fname)
                k = 1
                while os.path.exists(path):
                    fname = f"{base}({k}){ext}"
                    path = os.path.join(subj_dir, fname)
                    k += 1

                with open(path, "wb") as out:
                    out.write(raw)

                item = {
                    "name": fname,
                    "original_name": f.name,  # Сохраняем оригинальное название
                    "path": path,
                    "size": len(raw),
                    "mime": f.type or "application/octet-stream",
                    "uploader": user,
                    "note": note.strip(),
                    "ts": int(time.time())
                }
                items.append(item)
            save_materials_index(idx)
            st.success(f"Загружено файлов: {len(uploaded)}")

        # Управление файлами (только для преподавателей) - ПЕРЕД списком
        if current_role == "teacher" and items:
            with st.expander("🗑️ Управление файлами", expanded=False):
                st.markdown("**Выберите файлы для удаления:**")
                
                # Создаем чекбоксы для каждого файла
                files_to_delete = []
                for it in items:
                    display_name = display_filename(it['name'], it.get('original_name'))
                    file_info = f"{display_name} ({round(it['size']/1024,1)} KB)"
                    if it.get("note"):
                        file_info += f" - {it['note']}"
                    
                    if st.checkbox(file_info, key=f"del_check_{it['name']}"):
                        files_to_delete.append(it)
                
                if files_to_delete:
                    st.warning(f"Выбрано файлов для удаления: {len(files_to_delete)}")
                    if st.button("🗑️ Удалить выбранные файлы", type="primary"):
                        deleted_count = 0
                        for file_to_delete in files_to_delete:
                            try:
                                # Удаляем файл с диска
                                if os.path.exists(file_to_delete["path"]):
                                    os.remove(file_to_delete["path"])
                                # Удаляем из списка
                                items[:] = [it for it in items if it["name"] != file_to_delete["name"]]
                                deleted_count += 1
                            except Exception as e:
                                st.error(f"Ошибка удаления {file_to_delete['name']}: {str(e)}")
                        
                        if deleted_count > 0:
                            save_materials_index(idx)
                            st.success(f"✅ Удалено файлов: {deleted_count}")
                            st.rerun()
                else:
                    st.info("Выберите файлы для удаления")

        st.markdown("### Список материалов")
        if not items:
            st.info("Пока нет материалов для этого предмета.")
        else:
            # optional: filter by text
            q = st.text_input("Поиск по имени/описанию", placeholder="Например: лекция, дз2…")
            filtered = []
            if q:
                ql = q.lower()
                for it in items:
                    if ql in it["name"].lower() or (it.get("note") or "").lower().find(ql) >= 0:
                        filtered.append(it)
            else:
                filtered = items

            # Render list with download buttons
            for it in sorted(filtered, key=lambda x: x.get("ts", 0), reverse=True):
                with st.container(border=True):
                    display_name = display_filename(it['name'], it.get('original_name'))
                    st.markdown(f"**{display_name}**  —  <span class='muted'>{round(it['size']/1024,1)} KB</span>", unsafe_allow_html=True)
                    if it.get("note"):
                        st.markdown(f"<span class='badge'>Описание</span> {it['note']}", unsafe_allow_html=True)
                    st.markdown(f"<span class='muted'>Загрузил: {it.get('uploader','?')}</span>", unsafe_allow_html=True)

                    # Показываем аннотацию если есть
                    if it.get("annotation"):
                        with st.expander("📝 AI-аннотация", expanded=False):
                            st.markdown(it["annotation"])

                    # Кнопки действий
                    col_dl, col_ai = st.columns([1, 1])
                    
                    with col_dl:
                        # Read file for download button
                        try:
                            with open(it["path"], "rb") as fh:
                                data_bytes = fh.read()
                            st.download_button(
                                "⬇️ Скачать",
                                data=data_bytes,
                                file_name=it["name"],
                                mime=it.get("mime") or "application/octet-stream",
                                key=f"dl_{subj}_{it['name']}"
                            )
                        except FileNotFoundError:
                            st.error("Файл не найден на диске — возможно, был удалён вручную.")
                    
                    with col_ai:
                        # Кнопка генерации аннотации
                        if st.button("🤖 Создать аннотацию", key=f"ai_annot_{it['name']}"):
                            with st.spinner("AI анализирует материал..."):
                                # Извлекаем текст из файла
                                text = extract_text_from_file(it["path"], it.get("mime", ""))
                                
                                if text.startswith("Ошибка") or text.startswith("Формат файла"):
                                    st.error(text)
                                else:
                                    # Генерируем аннотацию
                                    annotation = generate_annotation(text, it["name"])
                                    
                                    # Сохраняем аннотацию в материал
                                    it["annotation"] = annotation
                                    save_materials_index(idx)
                                    
                                    st.success("✅ Аннотация создана!")
                                    st.rerun()

    
    with tab_videos:
        st.markdown("### 🎥 Видео материалы")
        
        # Функции для работы с видео
        def get_video_info(url):
            """Получает информацию о видео по ссылке"""
            try:
                if "youtube.com" in url or "youtu.be" in url:
                    # Извлекаем ID видео из YouTube ссылки
                    import re
                    video_id = None
                    if "youtube.com/watch?v=" in url:
                        video_id = url.split("v=")[1].split("&")[0]
                    elif "youtu.be/" in url:
                        video_id = url.split("youtu.be/")[1].split("?")[0]
                    
                    if video_id:
                        # Получаем название видео с YouTube
                        video_title = get_youtube_title(video_id)
                        
                        return {
                            "type": "youtube",
                            "video_id": video_id,
                            "embed_url": f"https://www.youtube.com/embed/{video_id}",
                            "title": video_title
                        }
                elif "vk.com" in url:
                    return {
                        "type": "vk",
                        "url": url,
                        "title": "VK видео"
                    }
                else:
                    return {
                        "type": "other",
                        "url": url,
                        "title": "Видео"
                    }
            except Exception as e:
                st.error(f"Ошибка обработки ссылки: {e}")
                return None
        
        def get_youtube_title(video_id):
            """Получает название YouTube видео по ID"""
            try:
                import requests
                from bs4 import BeautifulSoup
                
                # Создаем URL для получения страницы видео
                video_url = f"https://www.youtube.com/watch?v={video_id}"
                
                # Заголовки для имитации браузера
                headers = {
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
                }
                
                # Получаем страницу
                response = requests.get(video_url, headers=headers, timeout=10)
                response.raise_for_status()
                
                # Парсим HTML
                soup = BeautifulSoup(response.text, 'html.parser')
                
                # Ищем title в head
                title_tag = soup.find('title')
                if title_tag:
                    title = title_tag.get_text().strip()
                    # Убираем " - YouTube" из названия
                    if title.endswith(' - YouTube'):
                        title = title[:-10].strip()
                    return title
                
                # Если не нашли в title, ищем в meta тегах
                meta_title = soup.find('meta', property='og:title')
                if meta_title:
                    return meta_title.get('content', '').strip()
                
                # Если ничего не нашли, возвращаем ID
                return f"YouTube видео {video_id}"
                
            except Exception as e:
                # Если не удалось получить название, возвращаем ID
                return f"YouTube видео {video_id}"
        
        def load_videos():
            """Загружает список видео из файла"""
            try:
                with open("videos.json", "r", encoding="utf-8") as f:
                    return json.load(f)
            except FileNotFoundError:
                return {}
        
        def save_videos(videos):
            """Сохраняет список видео в файл"""
            with open("videos.json", "w", encoding="utf-8") as f:
                json.dump(videos, f, ensure_ascii=False, indent=2)
        
        # Загружаем видео для текущего предмета
        videos_data = load_videos()
        current_subject_videos = videos_data.get(current_subject(), [])
        
        # Форма добавления видео (только для преподавателей)
        if current_role == "teacher":
            st.markdown("#### Добавить видео")
            video_url = st.text_input("Ссылка на видео", placeholder="https://www.youtube.com/watch?v=...", key="video_url")
            video_note = st.text_input("Описание видео (опционально)", placeholder="Например: Лекция 1, Основы ML", key="video_note")
            
            if st.button("🎥 Добавить видео", use_container_width=True):
                if video_url and video_url.strip():
                    video_info = get_video_info(video_url.strip())
                    if video_info:
                        video_item = {
                            "url": video_url.strip(),
                            "title": video_info["title"],
                            "note": video_note.strip(),
                            "uploader": user,
                            "ts": int(time.time()),
                            "video_info": video_info
                        }
                        
                        current_subject_videos.append(video_item)
                        videos_data[current_subject()] = current_subject_videos
                        save_videos(videos_data)
                        
                        st.success(f"✅ Видео '{video_info['title']}' добавлено!")
                        st.rerun()
                    else:
                        st.error("❌ Не удалось обработать ссылку на видео")
                else:
                    st.error("❌ Введите ссылку на видео")
        else:
            st.info("🎥 Добавление видео доступно только преподавателям")
        
        # Отображение списка видео
        st.markdown("#### Список видео")
        if not current_subject_videos:
            st.info("Пока нет видео для этого предмета.")
        else:
            # Поиск по видео
            search_query = st.text_input("Поиск по названию/описанию", placeholder="Например: лекция, ML...", key="video_search")
            
            filtered_videos = current_subject_videos
            if search_query:
                search_lower = search_query.lower()
                filtered_videos = [v for v in current_subject_videos 
                                 if search_lower in v.get("title", "").lower() or 
                                    search_lower in v.get("note", "").lower()]
            
            # Отображаем видео
            for i, video in enumerate(sorted(filtered_videos, key=lambda x: x.get("ts", 0), reverse=True)):
                with st.container(border=True):
                    st.markdown(f"**{video['title']}**")
                    if video.get("note"):
                        st.markdown(f"<span class='badge'>Описание</span> {video['note']}", unsafe_allow_html=True)
                    st.markdown(f"<span class='muted'>Добавил: {video.get('uploader', '?')}</span>", unsafe_allow_html=True)
                    
                    # Встраиваем плеер
                    video_info = video.get("video_info", {})
                    if video_info.get("type") == "youtube":
                        st.components.v1.iframe(video_info["embed_url"], width=700, height=400)
                    else:
                        st.markdown(f"[Открыть видео]({video['url']})")
                    
                    # Кнопка удаления (только для преподавателей)
                    if current_role == "teacher":
                        if st.button(f"🗑️ Удалить", key=f"del_video_{i}"):
                            current_subject_videos.remove(video)
                            videos_data[current_subject()] = current_subject_videos
                            save_videos(videos_data)
                            st.success("Видео удалено")
                            st.rerun()

# Вкладка создания тестов (только для преподавателей)
if current_role == "teacher":
    with tab_test_creator:
        st.subheader("📝 Создание тестов")
        st.markdown("Создавайте тесты двух типов: с вариантами ответов или с оценкой по ключевым словам.")
        
        # Переключатель между созданием и управлением тестами
        tab_create, tab_ai_generate, tab_manage = st.tabs(["Создать тест", "🤖 AI-генерация", "Управление тестами"])
        
        with tab_create:
            st.markdown("### Создание нового теста")
            
            # Основная форма для создания теста 
            col1, col2 = st.columns(2)
            
            with col1:
                test_title = st.text_input("Название теста", placeholder="Например: Контрольная работа №1", key="test_title_input")
                test_subject = st.selectbox("Предмет", SUBJECTS, key="test_subject")
                # Срок сдачи (опционально)
                due_date = st.date_input("Срок сдачи (опционально)", key="test_due_date")
            
            with col2:
                test_type = st.selectbox(
                    "Тип теста",
                    ["multiple_choice", "keyword_based"],
                    format_func=lambda x: "С вариантами ответов" if x == "multiple_choice" else "С ключевыми словами",
                    key="test_type_selector"
                )
                test_description = st.text_area("Описание теста (опционально)", height=100, key="test_description_input")
                time_limit_minutes = st.number_input("Ограничение по времени (мин)", min_value=0, max_value=1440, value=0, help="0 — без ограничения", key="test_time_limit")
                
                # Сохраняем тип теста в session_state
                st.session_state.test_type = test_type
                
                # Автоматически обновляем форму при смене типа теста
                if "previous_test_type" not in st.session_state:
                    st.session_state.previous_test_type = test_type
                elif st.session_state.previous_test_type != test_type:
                    st.session_state.previous_test_type = test_type
                    st.session_state.questions = [{"question_id": "q1"}]
                    st.rerun()
            
            # Кнопка создания теста
            if st.button("🎯 Создать тест", use_container_width=True):
                # Собираем данные теста
                test_data = {
                    "title": test_title,
                    "subject": test_subject,
                    "test_type": test_type,
                    "description": test_description,
                    "questions": st.session_state.get("questions", [])
                }
                # Сохраняем срок сдачи, если установлен
                try:
                    if due_date:
                        # сериализуем как ISO-строку yyyy-mm-dd
                        test_data["due_date"] = due_date.isoformat()
                except Exception:
                    pass
                # Сохраняем таймер, если задан (>0)
                try:
                    if time_limit_minutes and int(time_limit_minutes) > 0:
                        test_data["time_limit_minutes"] = int(time_limit_minutes)
                except Exception:
                    pass
                
                # Генерируем ID теста
                test_id = generate_test_id(test_title, test_subject)
                test_data["assignment_id"] = test_id
                
                # Валидируем данные
                is_valid, error_message = validate_test_data(test_data)
                
                if is_valid:
                    # Добавляем тест в rubric.json
                    rubric_data = convert_to_rubric_format(test_data)
                    
                    # Загружаем существующие рубрики
                    with open("rubric.json", "r", encoding="utf-8") as f:
                        rubrics = json.load(f)
                    
                    # Добавляем новую рубрику
                    rubrics.append(rubric_data)
                    
                    # Сохраняем
                    with open("rubric.json", "w", encoding="utf-8") as f:
                        json.dump(rubrics, f, ensure_ascii=False, indent=2)
                    
                    st.success(f"✅ Тест '{test_title}' успешно создан и добавлен в задания!")
                    st.info(f"ID теста: {test_id}")
                    
                    # Очищаем форму
                    st.session_state.questions = [{"question_id": "q1"}]
                    st.rerun()
                else:
                    st.error(f"❌ Ошибка валидации: {error_message}")
            
            st.markdown("---")
            st.markdown("### Вопросы")
            
            # Динамическое добавление вопросов
            if "questions" not in st.session_state:
                st.session_state.questions = [{"question_id": "q1"}]
            
            # Кнопки управления вопросами (вне формы)
            col_add, col_clear = st.columns([1, 1])
            with col_add:
                if st.button("➕ Добавить вопрос", use_container_width=True):
                    new_q_id = f"q{len(st.session_state.questions) + 1}"
                    st.session_state.questions.append({"question_id": new_q_id})
                    st.rerun()
            
            with col_clear:
                if st.button("🗑️ Очистить все", use_container_width=True):
                    st.session_state.questions = [{"question_id": "q1"}]
                    st.rerun()
                
            # Инициализируем тип теста если его нет
            if "test_type" not in st.session_state:
                st.session_state.test_type = "multiple_choice"

            # Отображение вопросов
            for i, question in enumerate(st.session_state.questions):
                with st.expander(f"Вопрос {i+1} (ID: {question['question_id']})", expanded=True):
                    question_text = st.text_area(
                        f"Текст вопроса {i+1}",
                        value=question.get("question_text", ""),
                        key=f"q_text_{i}",
                        height=100
                    )
                    
                    # Получаем тип теста из session_state
                    test_type = st.session_state.get("test_type", "multiple_choice")
                    
                    if test_type == "multiple_choice":
                        st.markdown("**Варианты ответов:**")
                        
                        # Инициализируем options если их нет
                        if "options" not in st.session_state.questions[i]:
                            st.session_state.questions[i]["options"] = ["", ""]
                        
                        options = st.session_state.questions[i].get("options", ["", ""])
                        
                        # Отображаем существующие варианты
                        for j, option in enumerate(options):
                            option_text = st.text_input(
                                f"Вариант {j+1}",
                                value=option,
                                key=f"q_{i}_option_{j}"
                            )
                            # Обновляем вариант в session_state
                            st.session_state.questions[i]["options"][j] = option_text
                        
                        # Кнопки для добавления/удаления вариантов (вне формы)
                        col_buttons = st.columns([1, 1, 4])  # Две кнопки рядом, остальное место пустое
                        with col_buttons[0]:
                            if st.button("➕ Вариант", key=f"add_opt_{i}"):
                                st.session_state.questions[i]["options"].append("")
                                st.rerun()
                        with col_buttons[1]:
                            if st.button("➖ Вариант", key=f"rem_opt_{i}"):
                                if len(st.session_state.questions[i]["options"]) > 1:
                                    st.session_state.questions[i]["options"].pop()
                                    st.rerun()
                        
                        # Фильтруем пустые варианты для selectbox
                        valid_options = [opt for opt in st.session_state.questions[i]["options"] if opt.strip()]
                        
                        correct_answer = st.selectbox(
                            "Правильный ответ",
                            options=valid_options if valid_options else ["Нет вариантов"],
                            key=f"q_{i}_correct"
                        )
                        
                        max_points = st.number_input(
                            "Максимальные баллы",
                            min_value=1,
                            max_value=100,
                            value=question.get("max_points", 10),
                            key=f"q_{i}_points"
                        )
                        
                        # Обновляем данные вопроса
                        st.session_state.questions[i].update({
                            "question_text": question_text,
                            "correct_answer": correct_answer,
                            "max_points": max_points
                        })
                    
                    elif test_type == "keyword_based":
                        st.markdown("**Ключевые слова и баллы:**")
                        
                        # Инициализируем keywords если их нет
                        if "keywords" not in st.session_state.questions[i]:
                            st.session_state.questions[i]["keywords"] = [{"word": "", "points": 1}]
                        
                        keywords = st.session_state.questions[i].get("keywords", [{"word": "", "points": 1}])
                        
                        # Отображаем существующие ключевые слова
                        for j, keyword in enumerate(keywords):
                            col_word, col_points = st.columns([3, 1])
                            with col_word:
                                word = st.text_input(
                                    f"Ключевое слово {j+1}",
                                    value=keyword.get("word", ""),
                                    key=f"q_{i}_kw_{j}_word"
                                )
                            with col_points:
                                points = st.number_input(
                                    "Баллы",
                                    min_value=1,
                                    max_value=100,
                                    value=keyword.get("points", 1),
                                    key=f"q_{i}_kw_{j}_points"
                                )
                            
                            # Обновляем ключевое слово в session_state
                            st.session_state.questions[i]["keywords"][j] = {"word": word, "points": points}
                        
                        # Кнопки для добавления/удаления ключевых слов (вне формы)
                        col_buttons = st.columns([1, 1, 4])  # Две кнопки рядом, остальное место пустое
                        with col_buttons[0]:
                            if st.button("➕ Ключевое слово", key=f"add_kw_{i}"):
                                st.session_state.questions[i]["keywords"].append({"word": "", "points": 1})
                                st.rerun()
                        with col_buttons[1]:
                            if st.button("➖ Ключевое слово", key=f"rem_kw_{i}"):
                                if len(st.session_state.questions[i]["keywords"]) > 0:
                                    st.session_state.questions[i]["keywords"].pop()
                                    st.rerun()
                        
                        max_points = st.number_input(
                            "Максимальные баллы за вопрос",
                            min_value=1,
                            max_value=100,
                            value=question.get("max_points", 10),
                            key=f"q_{i}_max_points"
                        )
                        
                        # Обновляем данные вопроса
                        st.session_state.questions[i].update({
                            "question_text": question_text,
                            "max_points": max_points
                        })
        
        with tab_ai_generate:
            st.markdown("### 🤖 AI-генерация тестов")
            st.markdown("Создавайте тесты с помощью нейронной сети на основе доступных материалов.")
            
            # Проверяем доступность AI
            if client is None:
                st.error("❌ AI недоступен: не найден ключ API (DEEPSEEK_API_KEY)")
                st.info("Для использования AI-генерации тестов необходимо настроить API ключ.")
            else:
                # Форма для AI-генерации
                with st.form("ai_test_generation"):
                    col1, col2 = st.columns(2)
                    
                    with col1:
                        ai_test_title = st.text_input("Название теста", placeholder="Например: Контрольная по ML", key="ai_test_title")
                        ai_test_subject = st.selectbox("Предмет", SUBJECTS, key="ai_test_subject")
                        ai_test_description = st.text_area("Описание теста", placeholder="Опишите, что должен проверять тест", key="ai_test_description")
                    
                    with col2:
                        ai_question_count = st.number_input("Количество вопросов", min_value=1, max_value=20, value=5, key="ai_question_count")
                        ai_time_limit = st.number_input("Время на тест (минуты)", min_value=0, max_value=180, value=60, help="0 — без ограничения", key="ai_time_limit")
                        ai_due_date = st.date_input("Срок сдачи (опционально)", key="ai_due_date")
                    
                    # Выбор материалов для генерации
                    st.markdown("### 📚 Выберите материалы для генерации вопросов")
                    
                    # Получаем материалы для текущего предмета
                    idx, materials = get_materials_store()
                    if materials:
                        material_options = {}
                        for material in materials:
                            material_name = material.get("name", "")
                            material_note = material.get("note", "")
                            material_annotation = material.get("annotation", "")
                            
                            display_name = f"{display_filename(material_name, material.get('original_name'))}"
                            if material_note:
                                display_name += f" - {material_note}"
                            
                            material_options[display_name] = {
                                "name": material_name,
                                "note": material_note,
                                "annotation": material_annotation
                            }
                        
                        selected_materials = st.multiselect(
                            "Выберите материалы (можно несколько):",
                            options=list(material_options.keys()),
                            key="ai_selected_materials"
                        )
                    else:
                        st.warning("Нет доступных материалов для генерации вопросов.")
                        selected_materials = []
                    
                    # Кнопка генерации
                    if st.form_submit_button("🤖 Сгенерировать тест", use_container_width=True, type="primary"):
                        if ai_test_title and ai_test_description and selected_materials:
                            # Генерируем тест через AI
                            with st.spinner("AI генерирует тест..."):
                                generated_test = generate_ai_test(
                                    ai_test_title, 
                                    ai_test_subject, 
                                    ai_test_description,
                                    ai_question_count,
                                    selected_materials,
                                    material_options
                                )
                            
                            if generated_test:
                                # Сохраняем сгенерированный тест
                                if save_generated_test(generated_test, ai_time_limit, ai_due_date):
                                    st.success(f"✅ Тест '{ai_test_title}' успешно сгенерирован и сохранен!")
                                    st.rerun()
                                else:
                                    st.error("❌ Ошибка сохранения теста")
                            else:
                                st.error("❌ Ошибка генерации теста")
                        else:
                            st.error("❌ Заполните все обязательные поля и выберите материалы")
        
        with tab_manage:
            st.markdown("### Управление тестами")
            
            # Фильтр по предмету
            filter_subject = st.selectbox("Фильтр по предмету", ["Все"] + SUBJECTS, key="test_filter")
            
            # Получаем тесты из rubric.json (где они реально используются)
            with open("rubric.json", "r", encoding="utf-8") as f:
                all_rubrics = json.load(f)
            
            # Фильтруем тесты по предмету (показываем все тесты)
            if filter_subject == "Все":
                filtered_rubrics = all_rubrics
            else:
                filtered_rubrics = [r for r in all_rubrics if r.get("subject") == filter_subject]
            
            if not filtered_rubrics:
                st.info("Пока нет созданных тестов")
            else:
                # Отображение списка тестов
                for rubric in filtered_rubrics:
                    test_id = rubric["assignment_id"]
                    with st.container(border=True):
                        col_title, col_actions = st.columns([3, 1])
                        
                        with col_title:
                            st.markdown(f"**{rubric['title']}**")
                            
                            # Определяем тип теста
                            if rubric.get("test_type") == "multiple_choice":
                                test_type_display = "С вариантами ответов"
                            elif rubric.get("test_type") == "keyword_based":
                                test_type_display = "С ключевыми словами"
                            else:
                                # Старые тесты без test_type - определяем по структуре
                                if any("options" in q for q in rubric.get("questions", [])):
                                    test_type_display = "С вариантами ответов (старый формат)"
                                else:
                                    test_type_display = "С ключевыми словами (старый формат)"
                            
                            # Добавляем индикатор AI-генерации
                            ai_indicator = "🤖 AI-генерация" if rubric.get("ai_generated") else "✍️ Ручное создание"
                            
                            st.markdown(f"*Предмет: {rubric['subject']} | Тип: {test_type_display} | {ai_indicator}*")
                            st.markdown(f"*Вопросов: {len(rubric['questions'])} | ID: {test_id}*")
                        
                        with col_actions:
                            if st.button("🗑️", key=f"del_{test_id}", help="Удалить тест"):
                                # Удаляем из rubric.json
                                updated_rubrics = [r for r in all_rubrics if r["assignment_id"] != test_id]
                                
                                with open("rubric.json", "w", encoding="utf-8") as f:
                                    json.dump(updated_rubrics, f, ensure_ascii=False, indent=2)
                                
                                st.success("Тест удален")
                                st.rerun()
                            
                            if st.button("📋", key=f"edit_{test_id}", help="Просмотр и редактирование теста"):
                                # Переключаемся на режим редактирования
                                st.session_state[f"edit_test_{test_id}"] = True
                                st.rerun()
            
            # Проверяем, есть ли тест для редактирования
            editing_test = None
            for rubric in filtered_rubrics:
                test_id = rubric["assignment_id"]
                if st.session_state.get(f"edit_test_{test_id}", False):
                    editing_test = rubric
                    break
            
            # Отображение режима редактирования
            if editing_test:
                st.markdown("---")
                st.markdown(f"### ✏️ Редактирование теста: {editing_test['title']}")
                
                # Кнопка возврата
                if st.button("← Назад к списку", key="back_to_list"):
                    # Очищаем все флаги редактирования
                    for key in list(st.session_state.keys()):
                        if key.startswith("edit_test_"):
                            del st.session_state[key]
                    st.rerun()
                
                # Определяем тип теста
                test_type = editing_test.get("test_type")
                if not test_type:
                    if any("options" in q for q in editing_test.get("questions", [])):
                        test_type = "multiple_choice"
                    else:
                        test_type = "keyword_based"
                
                # Форма редактирования
                with st.form("edit_test_form"):
                    col1, col2 = st.columns(2)
                    
                    with col1:
                        new_title = st.text_input("Название теста", value=editing_test["title"])
                        new_subject = st.selectbox("Предмет", SUBJECTS, index=SUBJECTS.index(editing_test["subject"]) if editing_test["subject"] in SUBJECTS else 0)
                    
                    with col2:
                        new_test_type = st.selectbox(
                            "Тип теста",
                            ["multiple_choice", "keyword_based"],
                            index=0 if test_type == "multiple_choice" else 1,
                            format_func=lambda x: "С вариантами ответов" if x == "multiple_choice" else "С ключевыми словами"
                        )
                        new_description = st.text_area("Описание теста", value=editing_test.get("description", ""), height=100)
                    
                    st.markdown("---")
                    st.markdown("### Вопросы")
                    
                    # Отображение вопросов для редактирования
                    questions = editing_test.get("questions", [])
                    for i, question in enumerate(questions):
                        with st.expander(f"Вопрос {i+1} (ID: {question.get('question_id', f'q{i+1}')})", expanded=True):
                            question_text = st.text_area(
                                f"Текст вопроса {i+1}",
                                value=question.get("title", ""),
                                key=f"edit_q_text_{i}",
                                height=100
                            )
                            
                            if new_test_type == "multiple_choice":
                                st.markdown("**Варианты ответов:**")
                                options = question.get("options", [])
                                
                                for j, option in enumerate(options):
                                    option_text = st.text_input(
                                        f"Вариант {j+1}",
                                        value=option,
                                        key=f"edit_q_{i}_option_{j}"
                                    )
                                
                                correct_answer = st.selectbox(
                                    "Правильный ответ",
                                    options=options if options else ["Нет вариантов"],
                                    index=options.index(question.get("correct_answer", "")) if question.get("correct_answer", "") in options else 0,
                                    key=f"edit_q_{i}_correct"
                                )
                                
                                max_points = st.number_input(
                                    "Максимальные баллы",
                                    min_value=1,
                                    max_value=100,
                                    value=question.get("max_points", 10),
                                    key=f"edit_q_{i}_points"
                                )
                                
                            else:  # keyword_based
                                st.markdown("**Ключевые слова и баллы:**")
                                keywords = question.get("keywords", [])
                                
                                for j, keyword in enumerate(keywords):
                                    col_word, col_points = st.columns([3, 1])
                                    with col_word:
                                        word = st.text_input(
                                            f"Ключевое слово {j+1}",
                                            value=keyword.get("word", ""),
                                            key=f"edit_q_{i}_kw_{j}_word"
                                        )
                                    with col_points:
                                        points = st.number_input(
                                            "Баллы",
                                            min_value=1,
                                            max_value=100,
                                            value=keyword.get("points", 1),
                                            key=f"edit_q_{i}_kw_{j}_points"
                                        )
                                
                                max_points = st.number_input(
                                    "Максимальные баллы за вопрос",
                                    min_value=1,
                                    max_value=100,
                                    value=question.get("max_points", 10),
                                    key=f"edit_q_{i}_max_points"
                                )
                    
                    # Кнопки сохранения
                    col_save, col_cancel = st.columns([1, 1])
                    with col_save:
                        if st.form_submit_button("💾 Сохранить изменения", use_container_width=True):
                            # Собираем обновленные данные теста
                            updated_test = {
                                "subject": new_subject,
                                "assignment_id": editing_test["assignment_id"],
                                "title": new_title,
                                "test_type": new_test_type,
                                "description": new_description,
                                "questions": []
                            }
                            
                            # Обновляем вопросы
                            for i, question in enumerate(questions):
                                updated_question = {
                                    "question_id": question.get("question_id", f"q{i+1}"),
                                    "title": st.session_state.get(f"edit_q_text_{i}", question.get("title", "")),
                                    "max_points": question.get("max_points", 10)
                                }
                                
                                if new_test_type == "multiple_choice":
                                    # Обновляем варианты ответов
                                    options = []
                                    for j in range(len(question.get("options", []))):
                                        option_value = st.session_state.get(f"edit_q_{i}_option_{j}", "")
                                        if option_value:
                                            options.append(option_value)
                                    
                                    updated_question.update({
                                        "options": options,
                                        "correct_answer": st.session_state.get(f"edit_q_{i}_correct", ""),
                                        "max_points": st.session_state.get(f"edit_q_{i}_points", question.get("max_points", 10)),
                                        "test_type": "multiple_choice"
                                    })
                                else:
                                    # Обновляем ключевые слова
                                    keywords = []
                                    for j in range(len(question.get("keywords", []))):
                                        word = st.session_state.get(f"edit_q_{i}_kw_{j}_word", "")
                                        points = st.session_state.get(f"edit_q_{i}_kw_{j}_points", 1)
                                        if word:
                                            keywords.append({"word": word, "points": points})
                                    
                                    updated_question.update({
                                        "keywords": keywords,
                                        "max_points": st.session_state.get(f"edit_q_{i}_max_points", question.get("max_points", 10))
                                    })
                                
                                updated_test["questions"].append(updated_question)
                            
                            # Сохраняем изменения в rubric.json
                            with open("rubric.json", "r", encoding="utf-8") as f:
                                all_rubrics = json.load(f)
                            
                            # Находим и обновляем тест
                            for i, rubric in enumerate(all_rubrics):
                                if rubric["assignment_id"] == editing_test["assignment_id"]:
                                    all_rubrics[i] = updated_test
                                    break
                            
                            # Сохраняем обновленный файл
                            with open("rubric.json", "w", encoding="utf-8") as f:
                                json.dump(all_rubrics, f, ensure_ascii=False, indent=2)
                            
                            st.success("✅ Изменения сохранены!")
                            # Очищаем флаг редактирования
                            del st.session_state[f"edit_test_{editing_test['assignment_id']}"]
                            st.rerun()
                    
                    with col_cancel:
                        if st.form_submit_button("❌ Отмена", use_container_width=True):
                            # Очищаем флаг редактирования
                            del st.session_state[f"edit_test_{editing_test['assignment_id']}"]
                            st.rerun()
