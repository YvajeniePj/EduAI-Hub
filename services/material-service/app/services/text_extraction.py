"""
Text extraction service for various file formats
"""
import os
from typing import Optional


def extract_text_from_file(file_path: str, mime_type: str) -> str:
    """Extract text from file of different formats"""
    try:
        if mime_type == "application/pdf" or file_path.lower().endswith('.pdf'):
            try:
                import PyPDF2
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page in pdf_reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            # Try to fix encoding issues for Russian characters
                            try:
                                if 'Ð' in page_text or 'Ñ' in page_text:
                                    page_text = page_text.encode('latin1').decode('utf-8', errors='ignore')
                            except:
                                pass
                            text += page_text + "\n"
                    return text.strip()
            except Exception as pdf_error:
                return f"Error extracting text from PDF: {str(pdf_error)}"
        
        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or file_path.lower().endswith('.docx'):
            try:
                from docx import Document
                doc = Document(file_path)
                text = ""
                for paragraph in doc.paragraphs:
                    if paragraph.text:
                        text += paragraph.text + "\n"
                return text.strip()
            except Exception as docx_error:
                return f"Error extracting text from DOCX: {str(docx_error)}"
        
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or file_path.lower().endswith('.pptx'):
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text += shape.text + "\n"
                return text.strip()
            except Exception as pptx_error:
                return f"Error extracting text from PPTX: {str(pptx_error)}"
        
        elif mime_type == "text/plain" or file_path.lower().endswith('.txt'):
            # Try different encodings for TXT files
            encodings = ['utf-8', 'cp1251', 'latin1', 'utf-16']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        return file.read().strip()
                except UnicodeDecodeError:
                    continue
            return "Error: could not determine text file encoding"
        
        else:
            return f"File format {mime_type} is not supported for text extraction"
    
    except Exception as e:
        return f"Error extracting text: {str(e)}"

